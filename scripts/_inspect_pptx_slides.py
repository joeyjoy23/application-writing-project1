#!/usr/bin/env python
"""Inspect generated PPTX slide structure."""
from __future__ import annotations

import sys
from pathlib import Path

_PROJECT = Path(__file__).resolve().parents[1]
if str(_PROJECT) not in sys.path:
    sys.path.insert(0, str(_PROJECT))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = Path(r"D:\Downloads\ppt-work\mental_health_classroom.pptx")
OUT = Path(r"D:\Downloads\ppt-work\inspect-slides.txt")


def _shape_kind(sh) -> str:
    if sh.has_table:
        return f"table{len(sh.table.rows)}x{len(sh.table.columns)}"
    if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "auto"
    if sh.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        return "text"
    return str(sh.shape_type)


def main() -> None:
    prs = Presentation(str(PPTX))
    lines: list[str] = [f"pages={len(prs.slides)}", ""]
    issues: list[str] = []

    for si, slide in enumerate(prs.slides, 1):
        items: list[tuple[float, float, str, str]] = []
        for sh in slide.shapes:
            top = sh.top.inches if sh.top else 0.0
            h = sh.height.inches if sh.height else 0.0
            if sh.has_table:
                text = f"[TABLE {len(sh.table.rows)} rows]"
            elif sh.has_text_frame:
                text = (sh.text_frame.text or "").strip().replace("\n", " | ")
            else:
                continue
            if not text and not sh.has_table:
                continue
            items.append((top, h, _shape_kind(sh), text[:100]))
        items.sort(key=lambda x: x[0])
        bottom = max((t + h for t, h, _, _ in items), default=0.0)
        blank = max(0.0, (7.5 - bottom) / 7.5)
        title = items[0][3] if items else "(empty)"
        lines.append(f"=== Slide {si} blank={blank:.0%} bottom={bottom:.2f} shapes={len(items)} ===")
        lines.append(f"  title: {title}")
        for top, h, kind, text in items:
            lines.append(f"  {top:5.2f}+{h:.2f} [{kind:8}] {text}")

        # Heuristics
        if blank > 0.45 and len(items) <= 3:
            issues.append(f"S{si}: large_blank={blank:.0%} only {len(items)} text/table shapes")
        pills = [x for x in items if x[0] < 0.25 and x[1] <= 0.55 and "TABLE" not in x[3]]
        headers = [x for x in items if 0.72 <= x[0] <= 0.95 and x[1] <= 0.85 and len(x[3]) > 5]
        for p in pills:
            for hd in headers:
                if p[3] in hd[3] or hd[3].startswith(p[3][:4]):
                    issues.append(f"S{si}: redundant pill={p[3][:20]!r} header={hd[3][:30]!r}")
        if title == "句型" or (items and items[0][3] in ("句型", "词汇", "范文", "审题")):
            if len(items) <= 2 and blank > 0.4:
                issues.append(f"S{si}: section-only chrome blank={blank:.0%}")

    lines.extend(["", "=== ISSUES ==="])
    lines.extend(issues or ["(none)"])
    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {OUT} ({len(issues)} issues)")


if __name__ == "__main__":
    main()
