#!/usr/bin/env python
"""从应用文 HTML/Word 导出一键生成 V1 课堂 PPTX（python-pptx + WPS on-click）。"""

from __future__ import annotations

import argparse
import json
import logging
import sys
from pathlib import Path

from pptx import Presentation

_PROJECT_ROOT = Path(__file__).resolve().parents[1]
_SCRIPTS_DIR = _PROJECT_ROOT / "scripts"
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

from deck_plan import deck_plan_from_stage3, refine_deck_plan, replace_stage3_in_deck, stage3_specs_from_plan
from generate_classroom_pptx import (
    build_mental_health_deck,
    expand_slide_specs,
)
from wps_layout_verify import verify_deck_layout
from generate_classroom_pptx_v2 import _inject_v2_structure, render_v2_deck
from prepare_ppt_source import prepare_ppt_source
from ppt_work_cleanup import cleanup_ppt_work_dir
from pptx_click_reveal import apply_click_reveal, count_click_reveal_stats

log = logging.getLogger(__name__)


def _console_print(text: str) -> None:
    """Print safely on Windows consoles that lack full Unicode (e.g. GBK)."""
    try:
        print(text)
    except UnicodeEncodeError:
        enc = getattr(sys.stdout, "encoding", None) or "utf-8"
        print(text.encode(enc, errors="replace").decode(enc))


def _default_out_dir(export_path: Path) -> Path:
    return export_path.parent / "ppt-work"


_VOCAB_RETRY_LIMITS = (6, 4, 3, 2)


def _build_v2_slides(
    stage3_path: Path,
    deck_plan_path: Path | None,
    *,
    vocab_max_rows: int,
    use_custom_plan: bool,
    export_data_path: Path | None = None,
    preset: str = "70min",
    use_classroom_html: bool = False,
    module_dividers: bool = False,
    script_path: Path | None = None,
    template_id: str | None = None,
) -> list[dict]:
    from scripts.architecture_v1 import (
        build_architecture_deck,
        inject_module_dividers,
        load_export_data,
        merge_stage3_into_architecture_deck,
        parse_source_markdown_sections,
    )
    from scripts.classroom_script import (
        ClassroomScriptError,
        compile_classroom_script,
        load_classroom_script,
        load_script_template,
    )

    stage3_data = json.loads(stage3_path.read_text(encoding="utf-8"))
    if export_data_path and export_data_path.is_file():
        export_data = load_export_data(export_data_path)
    else:
        source_md = stage3_path.parent / "yingyongwen-source.md"
        export_data = (
            parse_source_markdown_sections(source_md)
            if source_md.is_file()
            else {}
        )

    lesson: str = preset if preset in ("40min", "70min", "80min") else "70min"
    if use_custom_plan and deck_plan_path and deck_plan_path.is_file():
        from deck_plan import load_deck_plan

        plan = load_deck_plan(deck_plan_path, stage3_data, vocab_max_rows=vocab_max_rows)
    else:
        from deck_plan import deck_plan_from_stage3, refine_deck_plan

        plan = deck_plan_from_stage3(stage3_data, vocab_max_rows=vocab_max_rows)
        plan = refine_deck_plan(stage3_data, plan, vocab_max_rows=vocab_max_rows)

    from deck_plan import stage3_specs_from_plan

    stage3_specs = stage3_specs_from_plan(stage3_data, plan)
    out_dir = stage3_path.parent
    if use_classroom_html:
        classroom_data_path = out_dir / "classroom-data.json"
        if classroom_data_path.is_file():
            from deck_from_classroom import classroom_to_deck

            classroom = json.loads(classroom_data_path.read_text(encoding="utf-8"))
            return classroom_to_deck(classroom, export_data, stage3_specs)
        print(
            "警告：--classroom-html 已指定但缺少 classroom-data.json，回退 Architecture V1",
            file=sys.stderr,
        )

    script: dict | None = None
    if script_path and script_path.is_file():
        script = load_classroom_script(script_path)
    elif template_id:
        script = load_script_template(template_id)
    else:
        default_script = out_dir / "classroom_script.json"
        if default_script.is_file():
            script = load_classroom_script(default_script)

    if script is not None:
        try:
            return compile_classroom_script(
                script,
                export_data,
                stage3_data,
                stage3_specs=stage3_specs,
                vocab_max_rows=vocab_max_rows,
            )
        except ClassroomScriptError as exc:
            print(f"警告：Classroom Script 编译失败（{exc}），回退 Architecture V1", file=sys.stderr)

    base = build_architecture_deck(export_data, preset=lesson)  # type: ignore[arg-type]
    merged = merge_stage3_into_architecture_deck(base, stage3_specs)
    return inject_module_dividers(merged, enabled=module_dividers)


def _render_and_verify(
    slides: list[dict],
    output_pptx: Path,
    *,
    no_anim: bool,
) -> tuple[dict, int]:
    from scripts.ppt_layout_fit import (
        expand_content_slides,
        expand_essay_slides,
        expand_peel_slides,
        expand_title_slides,
    )

    slides = expand_content_slides(
        expand_peel_slides(expand_essay_slides(expand_title_slides(slides)))
    )
    render_v2_deck(slides, output_pptx)
    if not no_anim:
        apply_click_reveal(output_pptx)
        timing_slides, anim_shapes = count_click_reveal_stats(output_pptx)
        log.info(
            "click-reveal applied: %d slides with timing, %d anim shapes",
            timing_slides,
            anim_shapes,
        )
    prs = Presentation(output_pptx)
    return verify_deck_layout(prs), len(prs.slides)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="从应用文导出一键生成 V1 课堂 PPTX（python-pptx + WPS on-click）"
    )
    parser.add_argument("export_path", type=Path, help="应用文 .html 或 .docx 导出路径")
    parser.add_argument(
        "-o",
        "--out-dir",
        type=Path,
        default=None,
        help="输出目录（默认：导出文件同级的 ppt-work）",
    )
    parser.add_argument(
        "--deck-plan",
        type=Path,
        default=None,
        help="classroom_deck.json（可选，覆盖默认 Stage3 页序）",
    )
    parser.add_argument(
        "--no-anim",
        action="store_true",
        help="跳过 WPS on-click 动画注入",
    )
    parser.add_argument(
        "--preset",
        choices=("40min", "70min", "80min"),
        default="70min",
        help="课时 preset：40min 精简 / 70min 标准（默认）/ 80min 完整（架构 V1）",
    )
    parser.add_argument(
        "--legacy-deck",
        action="store_true",
        help="使用旧版 mental_health 硬编码 deck（不推荐）",
    )
    parser.add_argument(
        "--v1",
        action="store_true",
        help="使用旧版 V1 扁平排版（无路线图/章节层次）",
    )
    parser.add_argument(
        "--keep-intermediate",
        action="store_true",
        help="保留 ppt-work 中间产物（json/md 等；默认生成后自动清理）",
    )
    parser.add_argument(
        "--classroom-html",
        action="store_true",
        help="旧版：用 -课件.html 的 16 页骨架映射 PPT（默认走 Architecture V1）",
    )
    parser.add_argument(
        "--module-dividers",
        action="store_true",
        help="插入 A–G 全屏模块分隔页（默认仅用页眉章节 tag）",
    )
    parser.add_argument(
        "--script",
        type=Path,
        default=None,
        help="classroom_script.json 路径（存在则走 Script 编译，否则 fallback Architecture V1）",
    )
    parser.add_argument(
        "--template",
        choices=("dual_poster_opinion", "letter_suggestion", "notice_campaign"),
        default=None,
        help="使用内置 Script 模板（无 --script 时生效）",
    )
    parser.add_argument(
        "--verbose",
        action="store_true",
        help="输出 WPS 布局校验调试信息",
    )
    args = parser.parse_args(argv)

    export_path = args.export_path.expanduser().resolve()
    if export_path.suffix.lower() not in {".html", ".docx"}:
        print("export_path 须为 .html 或 .docx", file=sys.stderr)
        return 1
    if not export_path.is_file():
        print(f"找不到导出文件：{export_path}", file=sys.stderr)
        return 1

    out_dir = (args.out_dir or _default_out_dir(export_path)).expanduser().resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    result = prepare_ppt_source(
        export_path,
        out_dir,
        write_classroom_data=args.classroom_html,
    )
    md_path, json_path = result[0], result[1]
    stage3_path = out_dir / "stage3.json"

    if not stage3_path.is_file():
        print(f"缺少 stage3.json：{stage3_path}", file=sys.stderr)
        return 1

    deck_plan = args.deck_plan.expanduser().resolve() if args.deck_plan else None
    if deck_plan and not deck_plan.is_file():
        fallback = out_dir / "classroom_deck.json"
        deck_plan = fallback if fallback.is_file() else None

    output_pptx = out_dir / "mental_health_classroom.pptx"
    export_data_path = out_dir / "export-data.json"
    layout_result: dict = {"ok": True, "pass1_issues": [], "pass2_issues": [], "wps_risk_count": 0}
    slide_count = 0

    if args.legacy_deck:
        from generate_classroom_pptx import build_mental_health_deck, expand_slide_specs
        from deck_plan import deck_plan_from_stage3, refine_deck_plan, replace_stage3_in_deck, stage3_specs_from_plan

        stage3_data = json.loads(stage3_path.read_text(encoding="utf-8"))
        plan = deck_plan_from_stage3(stage3_data)
        plan = refine_deck_plan(stage3_data, plan)
        stage3_specs = stage3_specs_from_plan(stage3_data, plan)
        slides = _inject_v2_structure(
            expand_slide_specs(replace_stage3_in_deck(build_mental_health_deck(), stage3_specs))
        )
        layout_result, slide_count = _render_and_verify(slides, output_pptx, no_anim=args.no_anim)
    elif args.v1:
        from generate_classroom_pptx import build_deck_with_stage3, render_deck

        slides = build_deck_with_stage3(stage3_path, deck_plan)
        render_deck(slides, output_pptx)
        if not args.no_anim:
            apply_click_reveal(output_pptx)
        prs = Presentation(output_pptx)
        layout_result = verify_deck_layout(prs)
        slide_count = len(prs.slides)
    else:
        use_custom = deck_plan is not None
        use_script = args.script is not None or args.template is not None
        retry_limits = (6,) if use_script else _VOCAB_RETRY_LIMITS
        for attempt, max_rows in enumerate(retry_limits):
            slides = _build_v2_slides(
                stage3_path,
                deck_plan,
                vocab_max_rows=max_rows,
                use_custom_plan=use_custom and attempt == 0,
                export_data_path=export_data_path,
                preset=args.preset,
                use_classroom_html=args.classroom_html,
                module_dividers=args.module_dividers,
                script_path=args.script.expanduser().resolve() if args.script else None,
                template_id=args.template,
            )
            layout_result, slide_count = _render_and_verify(
                slides, output_pptx, no_anim=args.no_anim
            )
            if layout_result["ok"]:
                if attempt > 0:
                    print(
                        f"Overflow retry succeeded on attempt {attempt + 1} "
                        f"(vocab_max_rows={max_rows})"
                    )
                break
            if attempt < len(retry_limits) - 1:
                n1 = len(layout_result["pass1_issues"])
                n2 = layout_result["wps_risk_count"]
                print(
                    f"Layout check failed (pass1={n1}, WPS_RISK={n2}); "
                    f"retry {attempt + 1} with vocab_max_rows={retry_limits[attempt + 1]}"
                )
                use_custom = False
            else:
                print(f"Overflow still failing after {len(retry_limits)} attempts", file=sys.stderr)

    print(f"架构: {'legacy' if args.legacy_deck else ('V1 flat' if args.v1 else f'Script/Architecture V1 ({args.preset})')}")
    print(f"导出文件: {export_path}")
    print(f"输出目录: {out_dir}")
    print(f"源稿: {md_path}")
    print(f"蓝图: {json_path}")
    print(f"Slides: {slide_count}")
    print(f"PPTX: {output_pptx}")
    print(f"动画: {'跳过' if args.no_anim else '已注入 on-click fade'}")
    wps_issues = layout_result["pass2_issues"]
    pass1_issues = layout_result["pass1_issues"]
    wps_report = layout_result.get("wps_report")
    print(
        f"Layout check: {'pass' if layout_result['ok'] else 'FAIL'} "
        f"(pass1={len(pass1_issues)}, WPS_RISK_OVERFLOW={layout_result['wps_risk_count']})"
    )
    if wps_report:
        print(
            f"WPS governance: critical={len(wps_report.critical_issues)} "
            f"warning={len(wps_report.warning_issues)} "
            f"cosmetic_ignored={len(wps_report.cosmetic_issues)} "
            f"teach_ready={wps_report.is_teach_ready} "
            f"risk_score={wps_report.risk_score:.3f}"
        )
    for issue in wps_issues[:5]:
        _console_print(f"  [CRITICAL] {issue}")
    warn = layout_result.get("warning_issues") or []
    for issue in warn[:3]:
        _console_print(f"  [WARNING]  {issue}")
    if len(wps_issues) > 5:
        print(f"  ... and {len(wps_issues) - 5} more critical issues")
    if args.verbose or wps_issues or warn:
        print("WPS layout governance report:")
        from wps_layout_verify import print_wps_layout_report

        print_wps_layout_report(Presentation(output_pptx))
    if pass1_issues and args.verbose:
        print("Pass1 (verify_text_fit) debug:")
        for issue in pass1_issues[:5]:
            _console_print(f"  [pass1] {issue}")

    if not args.keep_intermediate:
        removed = cleanup_ppt_work_dir(out_dir)
        if removed:
            print(f"已清理中间产物 {len(removed)} 项（保留对比用 PPTX）")
            for name in removed[:8]:
                print(f"  - {name}")
            if len(removed) > 8:
                print(f"  ... and {len(removed) - 8} more")

    return 0 if layout_result["ok"] else 1


if __name__ == "__main__":
    raise SystemExit(main())
