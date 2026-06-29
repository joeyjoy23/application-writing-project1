#!/usr/bin/env python
# DEPRECATED (2026-06-20): Do not use in default classroom PPT workflow.
# Regenerate SVGs via ppt-master Executor following PPT_LAYOUT_LAW instead.
# See: .cursor/skills/yingyongwen-export-to-ppt/PPT_LAYOUT_LAW.md
#      .cursor/skills/yingyongwen-export-to-ppt/SKILL.md
"""Fix V3 essay layout (wrap + badge width) and rebuild with ppt-master animations."""

from __future__ import annotations

import re
import shutil
import subprocess
import sys
import textwrap
from pathlib import Path

_PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

from scripts.generate_classroom_pptx import build_mental_health_deck

PPT_MASTER = Path(r"C:\Users\Joey\tools\ppt-master")
PROJECT = PPT_MASTER / "projects" / "yingyongwen-mental-health-v3_ppt169_20260620"
OUT_PPTX = Path(r"d:\Downloads\ppt-work\mental_health_classroom_V3_anim.pptx")

FONT_ESSAY = 35  # px in SVG → ~26pt in PPT after 0.75 conversion
LINE_H = 40
TEXT_X = 72
TEXT_MAX_W = 1140
WRAP_CHARS = 58
BODY_START_Y = 215
FOOTER_Y = 640


def _wrap_paragraph(text: str, width: int = WRAP_CHARS) -> list[str]:
    return textwrap.wrap(text, width=width, break_long_words=False, break_on_hyphens=False)


def _badge_path(width: int) -> str:
    w = max(200, min(360, width))
    return (
        f"M56,92 H{56 + w} A8,8 0 0 1 {64 + w},100 V120 A8,8 0 0 1 {56 + w},128 "
        f"H56 A8,8 0 0 1 48,120 V100 A8,8 0 0 1 56,92 Z"
    )


def build_essay_svg(
    title: str,
    badge_label: str,
    essay_text: str,
    footer: str,
) -> str:
    lines: list[tuple[str, int, str]] = []
    y = 175
    parts = [p.strip() for p in essay_text.split("\n\n") if p.strip()]
    if not parts:
        parts = [essay_text.strip()]

    for pi, para in enumerate(parts):
        wrapped = _wrap_paragraph(para)
        if not wrapped:
            continue
        if pi == 0 and wrapped[0].startswith("Dear"):
            lines.append((wrapped[0], y, "Times New Roman, Times, serif"))
            y += LINE_H
            wrapped = wrapped[1:]
        for line in wrapped:
            lines.append((line, y, "Times New Roman, Times, serif"))
            y += LINE_H
        y += 8

    if y > FOOTER_Y - 36:
        # tighten spacing if essay is long
        lines = []
        y = 175
        tight = 36
        for pi, para in enumerate(parts):
            wrapped = _wrap_paragraph(para, width=WRAP_CHARS + 4)
            if not wrapped:
                continue
            if pi == 0 and wrapped[0].startswith("Dear"):
                lines.append((wrapped[0], y, "Times New Roman, Times, serif"))
                y += tight
                wrapped = wrapped[1:]
            for line in wrapped:
                lines.append((line, y, "Times New Roman, Times, serif"))
                y += tight
            y += 6

    text_nodes = []
    reveal_groups = []
    idx = 1
    for line, ly, family in lines:
        t = (
            f'<text x="{TEXT_X}" y="{ly}" font-family="{family}" font-size="{FONT_ESSAY}" fill="#1E293B">{_esc(line)}</text>'
        )
        reveal_groups.append(f'  <g id="reveal-{idx:02d}">\n    {t}\n  </g>')
        idx += 1

    badge_w = int(len(badge_label) * 18 + 48)

    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720" width="1280" height="720">\n'
        f'  <rect width="1280" height="720" fill="#FFFFFF" />\n'
        f'  <g id="header">\n'
        f'    <rect x="0" y="0" width="1280" height="76" fill="#6366F1" />\n'
        f'    <rect x="0" y="76" width="1280" height="4" fill="#14B8A6" />\n'
        f'    <text x="48" y="52" font-family="Microsoft YaHei, Arial, sans-serif" font-size="51" font-weight="bold" fill="#FFFFFF">{_esc(title)}</text>\n'
        f'  </g>\n'
        f'  <g id="reveal-badge">\n'
        f'    <path fill="#14B8A6" d="{_badge_path(badge_w)}" />\n'
        f'    <text x="64" y="118" font-family="Microsoft YaHei, Arial, sans-serif" font-size="35" font-weight="bold" fill="#FFFFFF">{_esc(badge_label)}</text>\n'
        f'  </g>\n'
        f'  <g id="reveal-panel">\n'
        f'    <path fill="#F5F3FF" stroke="#E2E8F0" stroke-width="2" '
        f'd="M64,140 H1216 A16,16 0 0 1 1232,156 V644 A16,16 0 0 1 1216,660 H64 A16,16 0 0 1 48,644 V156 A16,16 0 0 1 64,140 Z" />\n'
        f'  </g>\n'
        f'{"".join(reveal_groups)}\n'
        f'  <g id="reveal-footer">\n'
        f'    <text x="{TEXT_X}" y="{FOOTER_Y}" font-family="Microsoft YaHei, Arial, sans-serif" font-size="35" fill="#64748B">{_esc(footer)}</text>\n'
        f'  </g>\n'
        f'</svg>\n'
    )


def _esc(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _essay_specs() -> list[tuple[str, str, str, str, str]]:
    deck = build_mental_health_deck()
    mapping = {
        "基础版范文 · 9分档": ("09_essay_basic", "110 words", "基础版 · 内容齐全、语言平实"),
        "高分版 A · 情感共鸣型": ("10_essay_advanced_a", "118 words", "12–14分档 · 情感共鸣 · 从句丰富"),
        "高分版 B · 逻辑思辨型": ("11_essay_advanced_b", "122 words", "12–14分档 · Firstly/Secondly · Unlike 比较"),
    }
    out: list[tuple[str, str, str, str, str]] = []
    for spec in deck:
        if spec.get("type") != "essay":
            continue
        title = spec["title"]
        if title not in mapping:
            continue
        stem, badge, footer = mapping[title]
        out.append((stem, title, badge, spec["essay_text"], footer))
    return out


def write_essay_svgs(target_dir: Path) -> None:
    target_dir.mkdir(parents=True, exist_ok=True)
    for stem, title, badge, essay, footer in _essay_specs():
        svg = build_essay_svg(title, badge, essay, footer)
        (target_dir / f"{stem}.svg").write_text(svg, encoding="utf-8")
        print(f"  wrote {stem}.svg ({len(svg)} bytes)")


def write_animations_json(project: Path) -> None:
    from scripts.ppt_v3_svg_patch import write_teaching_animations

    write_teaching_animations(project)


def sync_svg_dirs(project: Path) -> None:
    final_dir = project / "svg_final"
    out_dir = project / "svg_output"
    for src in final_dir.glob("*.svg"):
        shutil.copy2(src, out_dir / src.name)


def run_svg_to_pptx(project: Path) -> Path:
    py = PPT_MASTER / ".venv" / "Scripts" / "python.exe"
    script = PPT_MASTER / "skills" / "ppt-master" / "scripts" / "svg_to_pptx.py"
    cmd = [
        str(py),
        str(script),
        str(project),
        "-s",
        "final",
        "-t",
        "fade",
        "--transition-duration",
        "0.45",
        "-a",
        "fade",
        "--animation-duration",
        "0.5",
        "--animation-trigger",
        "on-click",
    ]
    print("  running svg_to_pptx with animations...")
    subprocess.run(cmd, check=True, cwd=str(PPT_MASTER))
    exports = sorted((project / "exports").glob("*.pptx"), key=lambda p: p.stat().st_mtime)
    if not exports:
        raise FileNotFoundError("No export pptx found")
    return exports[-1]


def postprocess_word_wrap(pptx_path: Path) -> None:
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE

    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell.text_frame.word_wrap = True
            if shape.has_text_frame:
                tf = shape.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.NONE
    prs.save(str(pptx_path))


def main() -> int:
    if not PROJECT.is_dir():
        print(f"Missing project: {PROJECT}")
        return 1

    print("Patching all SVG layout + animation groups...", flush=True)
    from scripts.ppt_v3_svg_patch import patch_all

    patch_all(PROJECT, write_essay_svgs)
    write_animations_json(PROJECT)

    print("Exporting PPTX with on-click animations...", flush=True)
    export = run_svg_to_pptx(PROJECT)
    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    dest = OUT_PPTX
    try:
        shutil.copy2(export, dest)
    except PermissionError:
        dest = OUT_PPTX.with_name("mental_health_classroom_V3_new.pptx")
        shutil.copy2(export, dest)
        print("  (target locked — saved as V3_new; close PowerPoint and rename)")
    postprocess_word_wrap(dest)
    print(f"V3 rebuilt: {dest}")
    print(f"  source export: {export}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
