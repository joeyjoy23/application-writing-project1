#!/usr/bin/env python
"""V2 classroom PPT — stronger hierarchy, youth appeal, logic-first layout."""

from __future__ import annotations

import re
import sys
from pathlib import Path

_PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from scripts.generate_classroom_pptx import (
    BOTTOM_Y,
    CONTENT_W,
    FONT_ANNOTATION,
    FONT_BADGE,
    FONT_BODY,
    FONT_ESSAY,
    FONT_ESSAY_SIZE,
    FONT_SECTION,
    FONT_TABLE,
    FONT_TITLE,
    FONT_UI,
    MARGIN_L,
    MIN_FONT_PT,
    SLIDE_H,
    SLIDE_W,
    WHITE,
    _clean_pptx_dir,
    _collect_font_sizes,
    _configure_text_frame,
    _enforce_word_wrap,
    _find_page_numbers,
    _scan_pptx_xml_pagination,
    build_deck_with_stage3,
    build_mental_health_deck,
    expand_slide_specs,
)
from scripts.wps_layout_verify import verify_deck_layout
from scripts.layout_solver_lite import (
    CARD_INNER_PAD_V,
    EssayPlan,
    SolvedLayout,
    TextBlockPlan,
    solve_content_layout,
    solve_fix_cards_layout,
    solve_layout,
    solve_table_row_heights,
    solve_title_layout,
)
from scripts.wps_layout_verify import compute_safe_height
from scripts.ppt_layout_fit import (
    ARROW_SEP_HEIGHT,
    LAYOUT_REGISTRY,
    MAX_CONTENT_BULLETS,
    PEEL_CARD_WIDTH,
    TABLE_MAX_ROWS_PER_PAGE,
    TEXT_BOX_PADDING_RATIO,
    _chunk_bullets_preserving_arrows,
    effective_text_area,
    essay_text_fits,
    fit_banner,
    fit_bullet_card_layout,
    fit_dual_cards,
    fit_essay_block,
    fit_paragraphs,
    fit_peel_point,
    fit_table_rows,
    fit_typography,
    is_arrow_separator,
    line_count,
    normalize_peel_point,
    peel_point_body_lines,
    pack_slides,
    split_banner_text,
    split_callout_lines,
    TITLE_PANEL_TOP,
    verify_safe_row_height,
    verify_safe_textbox_height,
    WPS_PANEL_FUDGE,
    text_block_height_paragraphs,
)

from styles.design_tokens import Theme as G
from styles.ppt_card import add_unified_card, inner_box_left, inner_box_width

INK = G.INK
MUTED = G.MUTED
WHITE = G.WHITE
BORDER = G.BORDER
WARN = G.WARNING
PANEL_SKY = G.ZONE_SURFACE
PANEL_LILAC = G.ZONE_SURFACE
PANEL_MINT = G.ZONE_SURFACE
PANEL_WARN = G.SURFACE
PRIMARY = G.PRIMARY
ACCENT = G.PRIMARY_DARK
HIGHLIGHT = G.ZONE_SURFACE
SECTION_COLORS = G.SECTION_COLORS

# Legacy module divider colors (architecture_v1)
CORAL = G.PRIMARY
INDIGO = G.PRIMARY_DARK
MINT = G.PRIMARY
SKY = G.ZONE_SURFACE
VIOLET = G.PRIMARY_DARK

CONTENT_TEXT_W = CONTENT_W

DEFAULT_V2_OUT = Path(r"d:\Downloads\ppt-work\mental_health_classroom.pptx")

_TIER_ACCENT = G.TIER_FILL


def _tier_color(label: str) -> RGBColor:
    return G.tier_style(label)[0]


def _tier_text_color(label: str) -> RGBColor:
    return G.tier_style(label)[1]


def _phrase_tier_note(tier: dict) -> str:
    level = tier.get("level", "")
    if level == "基础句":
        return ""
    note = tier.get("chinese", "") or ""
    hs = tier.get("high_score")
    if hs:
        note = f"{note}\n{hs}" if note else hs
    return note


def _inject_v2_structure(slides: list[dict]) -> list[dict]:
    """Mark inline section dividers on first slide of each arc (no standalone divider pages)."""
    out: list[dict] = []
    section_markers = {
        "审题 ·": ("01", "读懂题目", "写什么 · 写给谁 · 怎么写", G.DIVIDER_ACCENT),
        "PEEL": ("02", "范文骨架", "先搭 PEEL，再读三版范文", G.DIVIDER_ACCENT),
        "基础版范文": ("02", "范文骨架", "先搭 PEEL，再读三版范文", G.DIVIDER_ACCENT),
        "讲评活动": ("03", "讲评升级", "元素 → 主题 → 逻辑链", G.DIVIDER_ACCENT),
        "功能句型": ("04", "语言工具箱", "观点 · 论据 · 衔接 · 词块", G.DIVIDER_ACCENT),
        "当堂迁移": ("05", "当堂练", "用同样方法写新题", G.DIVIDER_ACCENT),
        "课堂小结": ("06", "带走什么", "审题 → 范文 → 句型 → 迁移", G.DIVIDER_ACCENT),
    }
    seen_sections: set[str] = set()

    for spec in slides:
        spec = dict(spec)
        title = spec.get("title", "")
        for key, (num, name, sub, color) in section_markers.items():
            if key in title and num not in seen_sections:
                spec["_inline_divider"] = {"num": num, "name": name, "subtitle": sub, "color": color}
                seen_sections.add(num)
                break
        out.append(spec)
    return out


class SlideBuilderV2:
    def __init__(self, prs: Presentation) -> None:
        self.prs = prs
        self._anim_seq = 0
        self._current_inline_divider: dict | None = None

    def _set_spec_context(self, spec: dict) -> None:
        self._current_inline_divider = spec.get("_inline_divider")

    def _reset_anim(self) -> None:
        self._anim_seq = 0

    def _tag_anim(self, shape) -> None:
        self._anim_seq += 1
        shape.name = f"anim_{self._anim_seq:03d}"

    def _unified_card(
        self,
        slide,
        left,
        top,
        width,
        height,
        *,
        warning: bool = False,
        accent: RGBColor | None = None,
    ):
        l = Inches(left) if isinstance(left, (int, float)) else left
        t = Inches(top) if isinstance(top, (int, float)) else top
        w = Inches(width) if isinstance(width, (int, float)) else width
        h = Inches(height) if isinstance(height, (int, float)) else height
        return add_unified_card(slide, l, t, w, h, accent=accent, warning=warning)

    def _pill(
        self,
        slide,
        text: str,
        *,
        left,
        top,
        width,
        height,
        fill: RGBColor,
        text_color: RGBColor = WHITE,
    ) -> None:
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left,
            top,
            width,
            height,
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = fill
        pill.line.fill.background()
        _configure_text_frame(pill.text_frame, anchor=MSO_ANCHOR.MIDDLE, pad_h=10, pad_v=6)
        p = pill.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT_UI
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

    def _blank(self):
        self._reset_anim()
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = G.BG
        bg.line.fill.background()
        return slide

    def _section_tag(self, slide, label: str) -> RGBColor:
        """Minimal section label — zone surface pill, no tier colors."""
        tag_w = min(11.5, 0.22 * len(label) + 1.2)
        tag = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.45),
            Inches(0.14),
            Inches(tag_w),
            Inches(0.40),
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = G.ZONE_SURFACE
        tag.line.color.rgb = G.BORDER
        tag.line.width = Pt(1)
        if tag.adjustments:
            tag.adjustments[0] = 0.12
        tf = tag.text_frame
        _configure_text_frame(tf, anchor=MSO_ANCHOR.MIDDLE, pad_h=6, pad_v=2)
        p = tf.paragraphs[0]
        p.text = label
        p.font.name = FONT_UI
        p.font.size = Pt(22)
        p.font.bold = False
        p.font.color.rgb = G.TEXT_SECONDARY
        p.alignment = PP_ALIGN.CENTER
        return G.PRIMARY

    def _simplify_title_for_section(self, title: str, section: str | None) -> str:
        if not section or not title:
            return title
        clean = title.strip()
        for sep in (" · ", " ·", "·"):
            prefix = f"{section}{sep}"
            if clean.startswith(prefix):
                return clean[len(prefix):].strip()
        if clean.startswith(section):
            return clean[len(section):].strip(" ·")
        return clean

    def _title_redundant_with_section(self, title: str, section: str | None) -> bool:
        if not section or not title:
            return False
        clean = title.strip()
        if clean.startswith(section):
            return True
        if f"{section} ·" in clean or f"{section}·" in clean:
            return True
        if section == "句型" and ("句型" in clean or "功能句型" in clean):
            return True
        if section == "词汇" and ("词块" in clean or "话题词块" in clean):
            return True
        if section == "范文" and "范文" in clean:
            return True
        if section == "审题" and ("审题" in clean or "思维" in clean):
            return True
        if section == "迁移" and "迁移" in clean:
            return True
        if section == "小结" and "小结" in clean:
            return True
        parts = [p.strip() for p in re.split(r"[·|]", clean) if p.strip()]
        return section in parts or any(p.startswith(section) for p in parts)

    def _inline_divider_strip(self, slide, marker: dict) -> float:
        """Thin module separator — zone band only."""
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(0),
            SLIDE_W,
            Inches(0.04),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = G.PRIMARY
        bar.line.fill.background()
        num = marker.get("num", "")
        name = marker.get("name", "")
        if num or name:
            lbl = slide.shapes.add_textbox(Inches(0.45), Inches(0.10), Inches(4.5), Inches(0.36))
            p = lbl.text_frame.paragraphs[0]
            p.text = f"{num}  {name}".strip()
            p.font.name = FONT_UI
            p.font.size = Pt(22)
            p.font.color.rgb = G.TEXT_SECONDARY
        return 0.52

    def _header(
        self,
        slide,
        title: str,
        section: str | None = None,
        *,
        inline_divider: dict | None = None,
    ) -> float:
        if inline_divider is None:
            inline_divider = self._current_inline_divider
        top_chrome = self._inline_divider_strip(slide, inline_divider) if inline_divider else 0.0
        redundant = self._title_redundant_with_section(title, section)
        compact = bool(section) and redundant

        if compact:
            self._section_tag(slide, section or title)
            bar_top = top_chrome + 0.52
            display = self._simplify_title_for_section(title, section) if section else title
            box = slide.shapes.add_textbox(
                Inches(0.55), Inches(bar_top), Inches(12.2), Inches(0.68)
            )
            tf = box.text_frame
            _configure_text_frame(tf, anchor=MSO_ANCHOR.MIDDLE)
            p = tf.paragraphs[0]
            p.text = display
            p.font.name = FONT_UI
            p.font.size = FONT_TITLE
            p.font.bold = True
            p.font.color.rgb = G.TEXT_TITLE
            return bar_top + 0.76

        if section:
            self._section_tag(slide, section)
        display_title = self._simplify_title_for_section(title, section) if section else title
        bar_top = top_chrome + 0.52
        box = slide.shapes.add_textbox(
            Inches(0.55), Inches(bar_top), Inches(12.2), Inches(0.68)
        )
        tf = box.text_frame
        _configure_text_frame(tf, anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.text = display_title
        p.font.name = FONT_UI
        p.font.size = FONT_TITLE
        p.font.bold = True
        p.font.color.rgb = G.TEXT_TITLE
        return bar_top + 0.76

    def title_slide(
        self,
        title: str,
        subtitle: str,
        body_lines: list[str],
        poster_lines: list[str] | None = None,
        *,
        task_tag: str = "",
    ) -> None:
        slide = self._blank()
        tbox = slide.shapes.add_textbox(MARGIN_L, Inches(0.34), Inches(7.8), Inches(0.74))
        _configure_text_frame(tbox.text_frame, anchor=MSO_ANCHOR.MIDDLE)
        p = tbox.text_frame.paragraphs[0]
        p.text = title
        p.font.name = FONT_UI
        p.font.size = FONT_TITLE
        p.font.bold = True
        p.font.color.rgb = G.TEXT_TITLE

        plan = solve_title_layout(
            {
                "type": "title",
                "title": title,
                "subtitle": subtitle,
                "body": body_lines,
                "poster_lines": poster_lines,
                "task_tag": task_tag or subtitle,
            }
        )
        if plan.task_pill_text:
            self._pill(
                slide,
                plan.task_pill_text,
                left=Inches(8.55),
                top=Inches(0.42),
                width=Inches(4.2),
                height=Inches(0.88),
                fill=G.TAG_BG,
                text_color=G.TEXT_SECONDARY,
            )

        if not plan.show_stem_panel and not plan.poster_lines:
            return

        panel_top_adj = TITLE_PANEL_TOP
        text_w = float(CONTENT_TEXT_W.inches) - 0.4
        pad_v = 0.18
        pad_h = 0.2

        if plan.show_stem_panel:
            panel_h = plan.stem_panel_height
            panel = self._unified_card(
                slide,
                float(MARGIN_L.inches),
                panel_top_adj,
                float(CONTENT_TEXT_W.inches),
                panel_h,
            )
            inner_h = max(0.45, panel_h - pad_v * 2)
            body = slide.shapes.add_textbox(
                MARGIN_L + Inches(pad_h + 0.06),
                Inches(panel_top_adj + pad_v),
                CONTENT_TEXT_W - Inches(pad_h * 2 + 0.06),
                Inches(inner_h),
            )
            _configure_text_frame(body.text_frame, pad_h=10, pad_v=8)
            tf = body.text_frame
            first = True
            for line in plan.stem_lines:
                para = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                para.text = line
                para.font.name = FONT_UI
                para.font.size = Pt(plan.stem_font_pt)
                para.font.color.rgb = INK
                para.space_after = Pt(10)
                para.line_spacing = 1.12
            self._tag_anim(panel)
            panel_top_adj += panel_h + 0.12

        if plan.poster_lines:
            poster_h = plan.poster_panel_height
            poster = self._unified_card(
                slide,
                float(MARGIN_L.inches),
                panel_top_adj,
                float(CONTENT_TEXT_W.inches),
                poster_h,
            )
            inner_h = max(0.45, poster_h - pad_v * 2)
            pbox = slide.shapes.add_textbox(
                MARGIN_L + Inches(pad_h + 0.06),
                Inches(panel_top_adj + pad_v),
                CONTENT_TEXT_W - Inches(pad_h * 2 + 0.06),
                Inches(inner_h),
            )
            _configure_text_frame(pbox.text_frame, pad_h=10, pad_v=8)
            tf_p = pbox.text_frame
            first_p = True
            for line in plan.poster_lines:
                para = tf_p.paragraphs[0] if first_p else tf_p.add_paragraph()
                first_p = False
                para.text = line
                para.font.name = FONT_UI
                para.font.size = Pt(plan.poster_font_pt)
                para.font.color.rgb = INK
                para.space_after = Pt(8)
                para.line_spacing = 1.12
            self._tag_anim(poster)

    def title_poster_slide(self, title: str, poster_lines: list[str]) -> None:
        """Second cover page when poster descriptions need their own slide."""
        slide = self._blank()
        top = self._header(slide, title, section="审题")
        poster_clean = [ln for ln in poster_lines if ln.strip()]
        if not poster_clean:
            return
        text_w = float(CONTENT_TEXT_W.inches) - 0.4
        avail_h = BOTTOM_Y - top - 0.12
        budget = LAYOUT_REGISTRY["title_body"]
        poster_fit = fit_paragraphs(
            poster_clean,
            text_w,
            avail_h - 0.36,
            space_after_pt=10,
            max_pt=budget.max_secondary_pt,
            min_pt=budget.min_pt,
            pad_h_pt=10,
            pad_v_pt=8,
        )
        panel_h = max(1.0, poster_fit.block_height * WPS_PANEL_FUDGE + 0.36)
        panel = self._unified_card(
            slide,
            float(MARGIN_L.inches),
            top + 0.08,
            float(CONTENT_TEXT_W.inches),
            panel_h,
        )
        pbox = slide.shapes.add_textbox(
            MARGIN_L + Inches(0.26),
            Inches(top + 0.26),
            CONTENT_TEXT_W - Inches(0.52),
            Inches(max(0.55, panel_h - 0.36)),
        )
        _configure_text_frame(pbox.text_frame, pad_h=10, pad_v=8)
        tf_p = pbox.text_frame
        first_p = True
        for line in poster_clean:
            para = tf_p.paragraphs[0] if first_p else tf_p.add_paragraph()
            first_p = False
            para.text = line
            para.font.name = FONT_UI
            para.font.size = Pt(poster_fit.font_pt)
            para.font.color.rgb = INK
            para.space_after = Pt(10)
            para.line_spacing = poster_fit.line_spacing

    def roadmap_slide(self, steps: list[tuple[str, str, str]]) -> None:
        slide = self._blank()
        top = self._header(slide, "本课怎么走？（逻辑路线图）", section="路线")
        col_w = 2.25
        gap = 0.15
        start_x = 0.55
        for i, (num, name, hint) in enumerate(steps):
            x = start_x + i * (col_w + gap)
            card = self._unified_card(slide, x, top + 0.15, col_w, 4.85)
            num_box = slide.shapes.add_textbox(
                Inches(x + 0.18),
                Inches(top + 0.32),
                Inches(0.55),
                Inches(0.55),
            )
            _configure_text_frame(num_box.text_frame, anchor=MSO_ANCHOR.MIDDLE)
            np = num_box.text_frame.paragraphs[0]
            np.text = num
            np.font.name = FONT_UI
            np.font.size = Pt(26)
            np.font.bold = True
            np.font.color.rgb = G.PRIMARY_DARK
            np.alignment = PP_ALIGN.CENTER

            tbox = slide.shapes.add_textbox(
                Inches(x + 0.12),
                Inches(top + 1.05),
                Inches(col_w - 0.24),
                Inches(3.2),
            )
            _configure_text_frame(tbox.text_frame, anchor=MSO_ANCHOR.TOP)
            p1 = tbox.text_frame.paragraphs[0]
            p1.text = name
            p1.font.name = FONT_UI
            p1.font.size = Pt(28)
            p1.font.bold = True
            p1.font.color.rgb = G.TEXT_TITLE
            p1.alignment = PP_ALIGN.CENTER
            p2 = tbox.text_frame.add_paragraph()
            p2.text = hint
            p2.font.name = FONT_UI
            p2.font.size = Pt(24)
            p2.font.color.rgb = G.TEXT_SECONDARY
            p2.alignment = PP_ALIGN.CENTER
            self._tag_anim(card)

            if i < len(steps) - 1:
                arrow = slide.shapes.add_textbox(
                    Inches(x + col_w - 0.05),
                    Inches(top + 2.05),
                    Inches(0.42),
                    Inches(0.55),
                )
                ap = arrow.text_frame.paragraphs[0]
                ap.text = "→"
                ap.font.name = FONT_UI
                ap.font.size = Pt(26)
                ap.font.color.rgb = G.TEXT_SECONDARY
                ap.alignment = PP_ALIGN.CENTER
                _configure_text_frame(arrow.text_frame, anchor=MSO_ANCHOR.MIDDLE)

        note = slide.shapes.add_textbox(MARGIN_L, Inches(6.55), CONTENT_W, Inches(0.55))
        p = note.text_frame.paragraphs[0]
        p.text = "【抓逻辑】每步都回答：我选哪个？为什么？用什么英语说？"
        p.font.name = FONT_UI
        p.font.size = Pt(24)
        p.font.color.rgb = G.TEXT_SECONDARY

    def divider_slide(self, num: str, name: str, subtitle: str, color: RGBColor) -> None:
        slide = self._blank()
        block = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
        )
        block.fill.solid()
        block.fill.fore_color.rgb = G.BG
        block.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(3.5), Inches(1.2))
        p = num_box.text_frame.paragraphs[0]
        p.text = num
        p.font.name = FONT_UI
        p.font.size = Pt(72)
        p.font.bold = True
        p.font.color.rgb = color

        name_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.5), Inches(1.0))
        p = name_box.text_frame.paragraphs[0]
        p.text = name
        p.font.name = FONT_UI
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = G.TEXT_TITLE

        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.05), Inches(11.0), Inches(0.8))
        p = sub_box.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = FONT_UI
        p.font.size = Pt(28)
        p.font.color.rgb = G.TEXT_SECONDARY

    def _guess_section(self, title: str) -> str | None:
        if "审题" in title or "动笔" in title or "易错" in title and "深化" not in title:
            return "审题"
        if "PEEL" in title or "范文" in title or "对比" in title or "讲评" in title:
            return "范文"
        if "词块" in title or "话题词" in title:
            return "词汇"
        if "句型" in title or "功能句型" in title:
            return "句型"
        if "迁移" in title:
            return "迁移"
        if "小结" in title:
            return "小结"
        if "活动" in title:
            return "活动"
        return None

    def _bullet_block_height(self, bullets: list[str], width: float) -> float:
        from scripts.wps_layout_verify import compute_card_height

        lines: list[str] = []
        for raw in bullets:
            display = (raw or "").strip()
            if is_arrow_separator(display):
                lines.append(display)
            elif display.startswith(("•", "→", "★", "①", "②")):
                lines.append(display)
            else:
                lines.append(f"• {display}")
        if not lines:
            return 0.0
        combined = "\n".join(lines)
        return compute_card_height(
            combined,
            28,
            max(3.0, width - 0.36),
            bullet=True,
            min_h=0.42,
        )

    def _single_bullet_card_height(self, bullet: str, width: float) -> float:
        from scripts.wps_layout_verify import compute_card_height

        raw = (bullet or "").strip()
        if is_arrow_separator(raw):
            return 0.32
        display = raw if raw.startswith("•") else f"• {raw}"
        body_h = compute_card_height(
            display,
            28,
            max(3.0, width - 0.52),
            bullet=True,
            min_h=0.42,
        )
        return body_h + 0.16

    def _paginate_bullets(self, bullets: list[str], avail_h: float, width: float) -> list[list[str]]:
        if not bullets:
            return [[]]
        pad = 0.18
        budget = max(0.5, avail_h - pad)
        pages: list[list[str]] = []
        current: list[str] = []
        for bullet in bullets:
            trial = current + [bullet]
            trial_h = self._bullet_block_height(trial, width)
            if current and trial_h > budget:
                pages.append(current)
                current = [bullet]
            else:
                current = trial
        if current:
            pages.append(current)
        if len(pages) == 2 and len(pages[-1]) <= 3:
            combined = pages[0] + pages[1]
            if self._bullet_block_height(combined, width) <= budget * 1.06:
                return [combined]
        return pages

    def _badge_redundant(self, title: str, badge: str | None) -> bool:
        if not badge:
            return True
        bt = badge.strip()
        tt = title.strip()
        if bt in tt or tt.startswith(bt) or tt.endswith(bt):
            return True
        if "迁移" in tt and "迁移" in bt:
            return True
        return False

    def _coalesce_bullet_pages(
        self, pages: list[list[str]], avail_h: float, width: float
    ) -> list[list[str]]:
        if len(pages) <= 1:
            return pages
        out = [list(p) for p in pages]
        while len(out) >= 2:
            combined = out[-2] + out[-1]
            if self._bullet_block_height(combined, width) <= avail_h:
                out = out[:-2] + [combined]
            elif len(out[-1]) <= 2 and self._bullet_block_height(combined, width) <= avail_h * 1.05:
                out = out[:-2] + [combined]
            else:
                break
        return out

    def _shift_essay_plan(self, plan: EssayPlan, header_bottom: float) -> EssayPlan:
        target_top = header_bottom + 0.08
        delta = target_top - plan.panel_top
        if abs(delta) < 0.02:
            return plan
        plan.panel_top = target_top
        if plan.annotation:
            plan.ann_top = target_top + plan.panel_height + 0.08
        return plan

    def _render_bullet_list(
        self,
        slide,
        bullets: list[str],
        top: float,
        left: float,
        width: float,
        avail_h: float,
        fill: RGBColor,
    ) -> None:
        """One unified card per page chunk (max 1 content block for bullets)."""
        del fill
        if not bullets:
            return
        pad_top = 0.14
        gap = 0.10
        block_h = min(avail_h, self._bullet_block_height(bullets, width) + pad_top * 2 + 0.08)
        warn = any("❌" in (b or "") or "别这样写" in (b or "") for b in bullets)
        card = self._unified_card(slide, left, top, width, block_h, warning=warn)
        inner_left = inner_box_left(left)
        inner_w = inner_box_width(width)
        y = top + pad_top
        for bullet in bullets:
            raw = bullet.strip()
            if is_arrow_separator(raw):
                ah = 0.30
                arrow_box = slide.shapes.add_textbox(
                    Inches(inner_left),
                    Inches(y),
                    Inches(inner_w),
                    Inches(ah),
                )
                _configure_text_frame(arrow_box.text_frame, anchor=MSO_ANCHOR.MIDDLE)
                ap = arrow_box.text_frame.paragraphs[0]
                ap.text = raw
                ap.font.name = FONT_UI
                ap.font.size = Pt(24)
                ap.font.color.rgb = G.TEXT_SECONDARY
                ap.alignment = PP_ALIGN.CENTER
                self._tag_anim(arrow_box)
                y += ah + gap
                continue

            display = raw if raw.startswith("•") else f"• {raw}"
            bullet_raw = display.lstrip("• ").strip()
            color = G.bullet_accent(display)
            bold = bullet_raw.startswith("★") or bullet_raw.startswith("❌")
            line_h = self._single_bullet_card_height(bullet, inner_w) - 0.08
            line_h = max(0.38, line_h)
            box = slide.shapes.add_textbox(
                Inches(inner_left),
                Inches(y),
                Inches(inner_w),
                Inches(line_h),
            )
            _configure_text_frame(box.text_frame, pad_h=2, pad_v=2, anchor=MSO_ANCHOR.TOP)
            para = box.text_frame.paragraphs[0]
            para.text = display
            para.font.name = FONT_UI
            para.font.size = Pt(28)
            para.font.bold = bold
            para.font.color.rgb = color
            para.line_spacing = 1.12
            para.space_after = Pt(4)
            self._tag_anim(box)
            y += line_h + gap

    def _content_section_title(self, slide, top: float, title: str) -> float:
        box = slide.shapes.add_textbox(MARGIN_L, Inches(top), CONTENT_W, Inches(0.56))
        _configure_text_frame(box.text_frame, anchor=MSO_ANCHOR.TOP)
        p = box.text_frame.paragraphs[0]
        p.text = title
        p.font.name = FONT_UI
        p.font.size = Pt(30)
        p.font.bold = True
        p.font.color.rgb = G.TEXT_TITLE
        self._tag_anim(box)
        return top + 0.62

    def content_slide(
        self,
        title: str,
        bullets: list[str],
        *,
        badge: str | None = None,
        panel: bool = False,
        warn_panel: bool = False,
        callout: str | None = None,
    ) -> None:
        slide = self._blank()
        section = self._guess_section(title)

        if callout and callout.strip():
            top = 0.36
            top = self._key_banner(slide, top, callout.strip(), label=None)
            section_title = title.split(" · ", 1)[-1].strip() if " · " in title else title
            top = self._content_section_title(slide, top + 0.14, section_title)
        else:
            top = self._header(slide, title, section=section)

        # Key insight banner (legacy bullets with 💡 prefix)
        key_lines = [b for b in bullets if b.startswith("💡") or "高分关键" in b or "最危险" in b]
        body_bullets = [b for b in bullets if b not in key_lines]
        if key_lines and not callout:
            budget = LAYOUT_REGISTRY["content_key"]
            banner_lines = split_banner_text(key_lines[0], budget)
            formula = " ".join(banner_lines).replace("抓重点", "").strip()
            if formula:
                body_bullets = [f"★ {formula}"] + body_bullets

        content_top = top + 0.08
        content_h = BOTTOM_Y - content_top
        if "PEEL" in title:
            self._peel_cards(slide, body_bullets, content_top, content_h)
            return

        list_w = float(CONTENT_W.inches) - 0.1
        list_left = float(MARGIN_L.inches) + 0.05
        pages = self._coalesce_bullet_pages(
            self._paginate_bullets(body_bullets, content_h, list_w),
            content_h,
            list_w,
        )
        for pi, chunk in enumerate(pages):
            if pi > 0:
                slide = self._blank()
                cont_title = f"{title}（续）"
                top = self._header(slide, cont_title, section=None)
                content_top = top + 0.08
                content_h = BOTTOM_Y - content_top
            self._render_bullet_list(
                slide, chunk, content_top, list_left, list_w, content_h, G.SURFACE
            )

    def _bullet_cards(
        self,
        slide,
        plans: list[TextBlockPlan],
        top: float,
        left: float,
        width: float,
        fill: RGBColor,
    ) -> None:
        if not plans:
            return
        tier_colors = {
            "基础": G.TEXT_MUTED,
            "进阶": G.TEXT_SECONDARY,
            "高级": G.TEXT_PRIMARY,
            "亮点": G.TEXT_PRIMARY,
            "必备": G.TEXT_MUTED,
        }
        gap = 0.14
        arrow_gap = 0.12
        y = top
        for plan in plans:
            if plan.container_height <= 0.14:
                arrow_box = slide.shapes.add_textbox(
                    Inches(left),
                    Inches(y),
                    Inches(width),
                    Inches(plan.container_height),
                )
                _configure_text_frame(arrow_box.text_frame, anchor=MSO_ANCHOR.MIDDLE)
                ap = arrow_box.text_frame.paragraphs[0]
                ap.text = plan.text.strip()
                ap.font.name = FONT_UI
                ap.font.size = Pt(28)
                ap.font.color.rgb = MUTED
                ap.alignment = PP_ALIGN.CENTER
                y += plan.container_height + arrow_gap
                continue

            card_h = plan.container_height
            inner_h = max(0.35, card_h - CARD_INNER_PAD_V * 2)
            accent = G.PRIMARY
            bullet_raw = plan.text.lstrip("• ").strip()
            if bullet_raw.startswith("★"):
                accent = G.SECONDARY
            elif bullet_raw.startswith("→"):
                accent = G.PRIMARY
            else:
                for tier, tc in tier_colors.items():
                    if bullet_raw.startswith(tier):
                        accent = tc
                        break

            card = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(y),
                Inches(width),
                Inches(card_h),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = fill
            card.line.color.rgb = accent
            card.line.width = Pt(2)

            bar = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(left),
                Inches(y + 0.06),
                Inches(0.08),
                Inches(max(0.2, card_h - 0.12)),
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent
            bar.line.fill.background()

            bold = plan.bold
            color = INK
            if bullet_raw.startswith("★"):
                color = G.SECONDARY
            elif bullet_raw.startswith("→"):
                color = G.PRIMARY
            else:
                for tier, tc in tier_colors.items():
                    if bullet_raw.startswith(tier):
                        color = tc
                        break

            box = slide.shapes.add_textbox(
                Inches(left + 0.22),
                Inches(y + CARD_INNER_PAD_V),
                Inches(width - 0.34),
                Inches(inner_h),
            )
            _configure_text_frame(box.text_frame, pad_h=6, pad_v=4)
            para = box.text_frame.paragraphs[0]
            para.text = plan.text
            para.font.name = FONT_UI
            para.font.size = Pt(plan.font_pt)
            para.font.bold = bold
            para.font.color.rgb = color
            para.line_spacing = plan.line_spacing
            para.space_after = Pt(5)
            self._tag_anim(card)
            y += card_h + gap

    def _peel_body_safe_height(self, point: dict, card_w: float) -> float:
        lines = peel_point_body_lines(point)
        text = "\n".join(lines)
        inner_w = max(3.0, card_w - 0.5)
        return compute_safe_height(text, 28, inner_w, bullet=True, min_h=0.85)

    def _peel_cards(self, slide, bullets: list[str], top: float, height: float) -> None:
        p1, p2 = [], []
        current = p1
        for b in bullets:
            if "Point 2" in b or "★ Point 2" in b:
                current = p2
            current.append(b)
        left_text = "\n".join(
            line.replace("核心句 P：", "P · ").replace("拓展 E：", "E · ")
            for line in p1
            if not line.startswith("★")
        )
        right_items = p2 or p1
        right_text = "\n".join(
            line.replace("核心句 P：", "P · ").replace("拓展 E：", "E · ")
            for line in right_items
            if not line.startswith("★")
        )
        w = 5.85
        left_h = compute_safe_height(left_text or " ", 28, w - 0.3, bullet=True, min_h=0.85)
        right_h = compute_safe_height(right_text or " ", 28, w - 0.3, bullet=True, min_h=0.85)
        card_h = max(left_h, right_h) + 0.72
        card_top = top
        groups = [
            ("① 选择", p1, left_h),
            ("② 理由", right_items, right_h),
        ]
        for idx, (label, items, body_h) in enumerate(groups):
            x = 0.55 + idx * (w + 0.25)
            card = self._unified_card(slide, x, card_top, w, card_h)
            hdr = slide.shapes.add_textbox(
                Inches(x + 0.22), Inches(card_top + 0.14), Inches(w - 0.36), Inches(0.48)
            )
            hp = hdr.text_frame.paragraphs[0]
            hp.text = label
            hp.font.name = FONT_UI
            hp.font.size = Pt(26)
            hp.font.bold = True
            hp.font.color.rgb = G.TEXT_TITLE
            inner_h = max(0.45, body_h)
            body = slide.shapes.add_textbox(
                Inches(x + 0.22),
                Inches(card_top + 0.58),
                Inches(w - 0.36),
                Inches(inner_h),
            )
            _configure_text_frame(body.text_frame)
            tf = body.text_frame
            for j, line in enumerate(items[:6]):
                if line.startswith("★"):
                    continue
                para = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                para.text = line.replace("核心句 P：", "P · ").replace("拓展 E：", "E · ")
                para.font.name = FONT_ESSAY if "I'd" in line or "The cracked" in line else FONT_UI
                para.font.size = Pt(28)
                para.font.color.rgb = INK
                para.space_after = Pt(5)
                para.line_spacing = 1.12
            self._tag_anim(body)

    def _render_essay_plan(self, slide, plan: EssayPlan) -> None:
        inner_h = max(0.45, plan.text_height)
        panel_h = inner_h + CARD_INNER_PAD_V * 2 + 0.12
        panel = self._unified_card(
            slide,
            float(MARGIN_L.inches),
            plan.panel_top - 0.06,
            float(CONTENT_W.inches),
            panel_h,
        )
        box = slide.shapes.add_textbox(
            MARGIN_L + Inches(0.22),
            Inches(plan.panel_top),
            CONTENT_W - Inches(0.36),
            Inches(inner_h),
        )
        _configure_text_frame(box.text_frame, anchor=MSO_ANCHOR.TOP, pad_h=10, pad_v=6)
        tf = box.text_frame
        first = True
        for pi, block in enumerate(plan.paragraphs):
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            para.text = (" " * plan.indent_spaces) + block.strip()
            para.font.name = FONT_ESSAY
            para.font.size = Pt(plan.font_pt)
            para.font.color.rgb = INK
            para.line_spacing = plan.line_spacing
            para.space_after = Pt(int(4 * 1.15))
            if pi > 0:
                para.space_before = Pt(
                    int((plan.para_space_pt if plan.line_spacing >= 1.0 else 4) * 1.15)
                )
        self._tag_anim(panel)

        if plan.annotation:
            ann_h = max(0.42, plan.ann_height + 0.16)
            ann_card = self._unified_card(
                slide,
                float(MARGIN_L.inches),
                plan.ann_top - 0.06,
                float(CONTENT_W.inches),
                ann_h,
            )
            note = slide.shapes.add_textbox(
                MARGIN_L + Inches(0.22),
                Inches(plan.ann_top),
                CONTENT_W - Inches(0.36),
                Inches(max(0.35, plan.ann_height)),
            )
            _configure_text_frame(note.text_frame, anchor=MSO_ANCHOR.TOP, pad_h=6, pad_v=4)
            np = note.text_frame.paragraphs[0]
            np.text = plan.annotation
            np.font.name = FONT_UI
            np.font.size = Pt(plan.ann_font_pt)
            np.font.color.rgb = MUTED
            np.line_spacing = plan.ann_line_spacing
            self._tag_anim(note)

    def essay_slide(
        self,
        title: str,
        essay_text: str,
        annotation: str,
        *,
        badge: str | None = None,
        layout_plan: EssayPlan | None = None,
    ) -> None:
        from scripts.essay_format import essay_layout_for_length, prepare_classroom_essay_display
        from scripts.layout_solver_lite import solve_essay_layout

        if layout_plan is not None:
            slide = self._blank()
            top = self._header(slide, title, section=None)
            self._render_essay_plan(slide, self._shift_essay_plan(layout_plan, top))
            return

        if "Dear " in essay_text or "中文批注" in essay_text:
            paragraphs, ann_text = prepare_classroom_essay_display(
                essay_text,
                annotation_fallback=annotation or "",
            )
        else:
            paragraphs = [p.strip() for p in essay_text.split("\n\n") if p.strip()]
            ann_text = (annotation or "").strip()
        line_spacing, para_space_pt, indent_spaces = essay_layout_for_length(
            [
                re.sub(r"\s*Word count:\s*\d+\s*$", "", p, flags=re.IGNORECASE).strip()
                for p in paragraphs
            ]
        )

        result = solve_essay_layout(
            paragraphs,
            annotation=ann_text,
            line_spacing=line_spacing,
            para_space_pt=para_space_pt,
            indent_spaces=indent_spaces,
        )
        if isinstance(result, list):
            for idx, part in enumerate(result):
                part_slide = self._blank()
                part_title = f"{title}（{idx + 1}/{len(result)}）" if len(result) > 1 else title
                top = self._header(part_slide, part_title, section=None)
                self._render_essay_plan(part_slide, self._shift_essay_plan(part, top))
            return

        slide = self._blank()
        top = self._header(slide, title, section=None)
        self._render_essay_plan(slide, self._shift_essay_plan(result, top))

    def table_slide(self, title: str, headers: list[str], rows: list[list[str]]) -> None:
        slide = self._blank()
        top = self._header(slide, title, section=None)
        n_cols = len(headers)
        col_fracs = [1.0 / n_cols] * n_cols
        kinds: list[str] = ["secondary"] * n_cols
        if n_cols:
            kinds[0] = "label"
        budget = LAYOUT_REGISTRY["content_cards"]
        max_table_h = BOTTOM_Y - top - 0.15
        per_row = min(0.72, max_table_h / max(len(rows) + 1, 1))
        row_heights = [0.56]
        row_fonts: list[list[int]] = [[26] * n_cols]
        row_spacings: list[list[float]] = [[1.12] * n_cols]
        for row in rows:
            rh, fpts, sps, _ = fit_table_rows(
                row, col_fracs, kinds, budget, per_row_budget=per_row, max_row_cap=per_row + 0.08
            )
            row_heights.append(rh)
            row_fonts.append(fpts)
            row_spacings.append(sps)
        total_h = sum(row_heights)
        if total_h > max_table_h:
            per_row = (max_table_h - 0.56) / max(len(rows), 1)
            row_heights = [0.56]
            for row in rows:
                rh, fpts, sps, _ = fit_table_rows(
                    row,
                    col_fracs,
                    kinds,
                    budget,
                    per_row_budget=max(0.38, per_row - 0.12),
                    max_row_cap=per_row,
                )
                row_heights.append(max(rh, 0.44))
            total_h = sum(row_heights)
        table_shape = slide.shapes.add_table(
            len(row_heights), n_cols, MARGIN_L, Inches(top + 0.12), CONTENT_W, Inches(total_h)
        )
        table = table_shape.table
        for ri, h in enumerate(row_heights):
            table.rows[ri].height = Inches(h)
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = G.TABLE_HEADER_BG
            _configure_text_frame(cell.text_frame, anchor=MSO_ANCHOR.MIDDLE)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_UI
                p.font.size = FONT_TABLE
                p.font.bold = True
                p.font.color.rgb = G.TABLE_HEADER
        for r, row in enumerate(rows, start=1):
            rh = row_heights[r]
            max_bh = 0.42
            for c, val in enumerate(row):
                col_w = float(CONTENT_W.inches) * col_fracs[c]
                kind = kinds[c]
                max_pt = budget.max_primary_pt if kind == "primary" else budget.max_secondary_pt
                if kind == "label":
                    max_pt = 26
                cell_fit = fit_typography(
                    val,
                    max(1.0, col_w - 0.12),
                    max(0.28, rh - 0.18),
                    min_pt=budget.min_pt,
                    max_pt=max_pt,
                )
                max_bh = max(max_bh, cell_fit.block_height)
                cell = table.cell(r, c)
                cell.text = val
                _configure_text_frame(cell.text_frame, anchor=MSO_ANCHOR.MIDDLE)
                for p in cell.text_frame.paragraphs:
                    p.font.name = FONT_UI
                    p.font.size = Pt(cell_fit.font_pt)
                    p.font.color.rgb = INK if c == 0 else MUTED
                    p.line_spacing = cell_fit.line_spacing
                    if c > 0:
                        p.font.bold = False
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = G.ZONE_SURFACE
            row_heights[r] = max(rh, max_bh + 0.24)
            table.rows[r].height = Inches(row_heights[r])
        table_shape.height = Inches(sum(row_heights))

    def _key_banner(self, slide, top: float, text: str, *, label: str | None = "抓重点") -> float:
        if label is None:
            text_w = float(CONTENT_W.inches) - 0.44
            lines, banner_fit, h = split_callout_lines(text, text_w)
        else:
            budget = LAYOUT_REGISTRY["phrase_table_footer"]
            lines = split_banner_text(text, budget)
            banner_fit = fit_banner("\n".join(lines), budget)
            h = max(0.92, banner_fit.block_height + 0.22)
            if len(lines) > 1:
                h = max(1.08, banner_fit.block_height + 0.32)
        card = self._unified_card(
            slide,
            float(MARGIN_L.inches),
            top,
            float(CONTENT_W.inches),
            h,
        )
        body_left = MARGIN_L + Inches(0.26)
        body_width = Inches(CONTENT_W.inches - 0.52)
        if label:
            lbl = slide.shapes.add_textbox(MARGIN_L + Inches(0.26), Inches(top + 0.10), Inches(1.4), Inches(0.44))
            lp = lbl.text_frame.paragraphs[0]
            lp.text = label
            lp.font.name = FONT_UI
            lp.font.size = Pt(22)
            lp.font.color.rgb = G.TEXT_SECONDARY
            body_left = MARGIN_L + Inches(1.72)
            body_width = Inches(CONTENT_W.inches - 1.98)
        body = slide.shapes.add_textbox(
            body_left, Inches(top + 0.12), body_width, Inches(h - 0.2)
        )
        _configure_text_frame(body.text_frame)
        first = True
        for line in lines:
            bp = body.text_frame.paragraphs[0] if first else body.text_frame.add_paragraph()
            first = False
            bp.text = line
            bp.font.name = FONT_UI
            bp.font.size = Pt(banner_fit.font_pt)
            bp.font.bold = label is None
            bp.font.color.rgb = G.TEXT_BODY
            bp.line_spacing = banner_fit.line_spacing
        self._tag_anim(card)
        gap = 0.20 if label is None else 0.12
        return top + h + gap

    def _style_table_header(self, table, headers: list[str]) -> None:
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = G.TABLE_HEADER_BG
            _configure_text_frame(cell.text_frame, anchor=MSO_ANCHOR.MIDDLE)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_UI
                p.font.size = FONT_TABLE
                p.font.bold = True
                p.font.color.rgb = G.TABLE_HEADER

    def _style_phrase_row(
        self,
        table,
        row_idx: int,
        level: str,
        english: str,
        note: str,
        font_pts: list[int] | None = None,
        line_spacings: list[float] | None = None,
    ) -> None:
        accent = _tier_color(level)
        accent_text = _tier_text_color(level)
        vals = [level, english, note]
        sizes = font_pts or [26, 28, 26]
        spacings = line_spacings or [1.12, 1.12, 1.12]
        for c, val in enumerate(vals):
            cell = table.cell(row_idx, c)
            cell.text = val
            anchor = MSO_ANCHOR.MIDDLE if c == 0 else MSO_ANCHOR.TOP
            _configure_text_frame(cell.text_frame, anchor=anchor)
            pt = sizes[c] if c < len(sizes) else 26
            sp = spacings[c] if c < len(spacings) else 1.12
            for p in cell.text_frame.paragraphs:
                if c == 0:
                    p.font.name = FONT_UI
                    p.font.bold = True
                    p.font.color.rgb = accent_text
                elif c == 1:
                    p.font.name = FONT_ESSAY
                    p.font.bold = True
                    p.font.color.rgb = INK
                else:
                    p.font.name = FONT_UI
                    p.font.color.rgb = MUTED
                p.font.size = Pt(pt)
                p.line_spacing = sp
            if c == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = accent
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL_SKY

    def _render_peel_point_body(
        self,
        slide,
        *,
        x: float,
        card_top: float,
        card_w: float,
        body_h: float,
        point: dict,
        color: RGBColor,
        body_fit,
    ) -> None:
        body = slide.shapes.add_textbox(
            Inches(x + 0.15),
            Inches(card_top + 0.55),
            Inches(card_w - 0.3),
            Inches(body_h),
        )
        _configure_text_frame(body.text_frame, pad_h=4, pad_v=4)
        tf = body.text_frame
        first = True
        norm = normalize_peel_point(point)

        def add_line(tag: str, text: str, *, essay: bool = False, accent: RGBColor = INK) -> None:
            nonlocal first
            if not text:
                return
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            run_tag = para.add_run()
            run_tag.text = f"{tag}  "
            run_tag.font.name = FONT_UI
            run_tag.font.size = Pt(26)
            run_tag.font.bold = True
            run_tag.font.color.rgb = color
            run_body = para.add_run()
            run_body.text = text
            run_body.font.name = FONT_ESSAY if essay else FONT_UI
            run_body.font.size = Pt(body_fit.font_pt)
            run_body.font.color.rgb = accent
            para.space_after = Pt(4 if body_fit.line_spacing < 1.0 else 6)
            para.line_spacing = body_fit.line_spacing

        add_line("P", norm["p"], essay=True, accent=INK)
        for ei, e_text in enumerate(norm["e_items"]):
            add_line(
                "E" if ei == 0 else "  ",
                e_text,
                essay="'" in e_text or "I'd" in e_text,
            )
        add_line("L", norm["l"], essay=True, accent=G.PRIMARY)
        self._tag_anim(body)

    def _peel_card(
        self,
        slide,
        *,
        x: float,
        card_top: float,
        card_w: float,
        card_h: float,
        label: str,
        point: dict,
        color: RGBColor,
        fill: RGBColor,
        body_fit,
    ) -> None:
        lines = peel_point_body_lines(point)
        inner_w = max(3.0, card_w - 0.5)
        body_h = compute_safe_height(
            "\n".join(lines),
            body_fit.font_pt,
            inner_w,
            bullet=True,
            min_h=0.85,
        )
        card_h = max(card_h, body_h + 0.72)

        card = self._unified_card(slide, x, card_top, card_w, card_h)
        hdr = slide.shapes.add_textbox(
            Inches(x + 0.22), Inches(card_top + 0.12), Inches(card_w - 0.36), Inches(0.45)
        )
        hp = hdr.text_frame.paragraphs[0]
        hp.text = (point.get("heading") or label).replace("★", "").strip()
        hp.font.name = FONT_UI
        hp.font.size = Pt(26)
        hp.font.bold = True
        hp.font.color.rgb = G.TEXT_TITLE

        self._render_peel_point_body(
            slide,
            x=x + 0.06,
            card_top=card_top,
            card_w=card_w - 0.12,
            body_h=body_h,
            point=point,
            color=G.TEXT_BODY,
            body_fit=body_fit,
        )

    def peel_slide(self, title: str, points: list[dict], *, layout: str = "dual") -> None:
        points = [normalize_peel_point(p) for p in points[:2]]
        if not points:
            return
        if layout == "single" or len(points) == 1:
            slide = self._blank()
            top = self._header(slide, title, section=None)
            card_top = top + 0.08
            card_w = CONTENT_W.inches if hasattr(CONTENT_W, "inches") else 12.0
            body_fit = fit_peel_point(points[0], card_w, 99.0)
            body_h = self._peel_body_safe_height(points[0], card_w)
            card_h = body_h + 0.72
            self._peel_card(
                slide,
                x=0.55,
                card_top=card_top,
                card_w=card_w,
                card_h=card_h,
                label="① 先选" if "1" in points[0].get("label", "1") else "② 再讲理由",
                point=points[0],
                color=G.ACCENT if "1" in points[0].get("label", "1") else G.TEXT_PRIMARY,
                fill=G.SURFACE,
                body_fit=body_fit,
            )
            return

        slide = self._blank()
        top = self._header(slide, title, section=None)
        card_top = top + 0.08
        card_w = PEEL_CARD_WIDTH
        body_heights = [self._peel_body_safe_height(p, card_w) for p in points]
        fits = [fit_peel_point(p, card_w, 99.0) for p in points]
        card_h = max(body_heights) + 0.72
        labels = ("① 先选", "② 再讲理由")
        colors = (G.ACCENT, G.TEXT_PRIMARY)
        fills = (G.SURFACE, G.SURFACE)
        for idx, point in enumerate(points):
            x = 0.55 + idx * (card_w + 0.25)
            self._peel_card(
                slide,
                x=x,
                card_top=card_top,
                card_w=card_w,
                card_h=card_h,
                label=labels[idx],
                point=point,
                color=colors[idx],
                fill=fills[idx],
                body_fit=fits[idx],
            )

    def phrase_table_slide(
        self,
        title: str,
        table: dict,
        *,
        badge: str | None = None,
        part: str = "full",
    ) -> None:
        """part: body | footer | full — body and footer always render as separate slides."""
        if part in ("full", "body"):
            self._phrase_table_body_slide(title, table, badge=badge)
        if part in ("full", "footer", "footer_note", "footer_fix", "footer_fix_bad", "footer_fix_good"):
            if part == "full" and not (
                table.get("topic_note") or table.get("fix_bad") or table.get("fix_good")
            ):
                return
            footer_title = (
                title
                if part.startswith("footer")
                else f"{title.split('·')[0].strip()} · 用法与改错"
            )
            self._phrase_table_footer_slide(
                footer_title,
                table,
                part=part if part.startswith("footer") else "footer",
            )

    def _phrase_table_body_slide(
        self, title: str, table: dict, *, badge: str | None = None
    ) -> None:
        slide = self._blank()
        top = self._header(slide, title, section="句型")

        col_fracs = [0.14, 0.52, 0.34]
        tiers = table.get("tiers", [])
        content_top = top + 0.06
        row_data: list[tuple[str, str, str]] = []
        for tier in tiers:
            row_data.append(
                (tier.get("level", ""), tier.get("english", ""), _phrase_tier_note(tier))
            )

        row_heights = solve_table_row_heights(
            [[level, english, note] for level, english, note in row_data],
            col_fracs,
            header=True,
        )
        row_fonts = [[26, 26, 26]]
        row_spacings = [[1.12, 1.12, 1.12]]
        budget = LAYOUT_REGISTRY["phrase_table_body"]
        for ri, (level, english, note) in enumerate(row_data, start=1):
            rh = row_heights[ri]
            fpts = [26, 28, 26]
            fsps = [1.12, 1.12, 1.12]
            row_fonts.append(fpts)
            row_spacings.append(fsps)

        tbl_shape = slide.shapes.add_table(
            len(row_heights), 3, MARGIN_L, Inches(content_top), CONTENT_W, Inches(sum(row_heights))
        )
        tbl = tbl_shape.table
        for i, frac in enumerate(col_fracs):
            tbl.columns[i].width = Inches(CONTENT_W.inches * frac)
        for ri, h in enumerate(row_heights):
            tbl.rows[ri].height = Inches(h)
        self._style_table_header(tbl, ["层级", "背这句", "怎么用"])
        for ri, (level, english, note) in enumerate(row_data, start=1):
            rh = row_heights[ri]
            fpts = row_fonts[ri] if ri < len(row_fonts) else [26, 28, 26]
            fsps = row_spacings[ri] if ri < len(row_spacings) else [1.12, 1.12, 1.12]
            final_pts = fpts
            final_sps = fsps
            self._style_phrase_row(
                tbl, ri, level, english, note, font_pts=final_pts, line_spacings=final_sps
            )
        tbl_shape.height = Inches(sum(row_heights))

    def _phrase_table_footer_slide(self, title: str, table: dict, *, part: str = "footer") -> None:
        slide = self._blank()
        top = self._header(slide, title, section="句型")
        cursor = top + 0.12
        budget = LAYOUT_REGISTRY["phrase_table_footer"]
        show_note = part in ("footer", "footer_note") and table.get("topic_note")
        show_fix = part in ("footer", "footer_fix", "footer_fix_bad", "footer_fix_good") and (
            table.get("fix_bad") or table.get("fix_good")
        )
        if show_note:
            cursor = self._key_banner(slide, cursor, table["topic_note"], label="本题")
            for sh in slide.shapes:
                if sh.has_text_frame and table["topic_note"][:12] in (sh.text_frame.text or ""):
                    self._tag_anim(sh)
                    break

        if not show_fix:
            return

        col_w = 5.85
        y = cursor + 0.08
        avail_h = BOTTOM_Y - y - 0.12
        left_text = ("别这样写\n" + table.get("fix_bad", "")).strip() if table.get("fix_bad") else ""
        right_text = (
            ("改成\n" + table.get("fix_good", "").lstrip("→").strip()).strip()
            if table.get("fix_good")
            else ""
        )
        if part == "footer_fix_bad":
            left_text = left_text or ""
            right_text = ""
        elif part == "footer_fix_good":
            left_text = ""
            right_text = right_text or ""

        fix_mode, fix_plans = solve_fix_cards_layout(
            left_text,
            right_text,
            avail_height=max(0.5, avail_h),
            card_width=col_w,
            full_width=float(CONTENT_W.inches),
        )

        def _draw_fix_card_plan(
            x: float,
            card_y: float,
            card_w: float,
            plan,
            *,
            warning: bool = False,
            bold: bool = False,
        ) -> None:
            card_h = plan.container_height
            inner_h = max(0.35, card_h - CARD_INNER_PAD_V * 2)
            card = self._unified_card(slide, x, card_y, card_w, card_h, warning=warning)
            box = slide.shapes.add_textbox(
                Inches(x + 0.22),
                Inches(card_y + CARD_INNER_PAD_V),
                Inches(card_w - 0.36),
                Inches(inner_h),
            )
            _configure_text_frame(box.text_frame, pad_h=6, pad_v=4)
            bp = box.text_frame.paragraphs[0]
            bp.text = plan.text
            bp.font.name = FONT_ESSAY
            bp.font.size = Pt(plan.font_pt)
            bp.font.bold = bold
            bp.font.color.rgb = G.TEXT_BODY
            bp.line_spacing = plan.line_spacing
            self._tag_anim(box)

        if fix_mode == "stack" and len(fix_plans) >= 2:
            cy = y
            for idx, plan in enumerate(fix_plans):
                is_bad = "别这样写" in plan.text
                _draw_fix_card_plan(
                    float(MARGIN_L.inches),
                    cy,
                    float(CONTENT_W.inches),
                    plan,
                    warning=is_bad,
                    bold=not is_bad,
                )
                cy += plan.container_height + 0.08
        elif fix_mode == "dual" and len(fix_plans) >= 2:
            for idx, plan in enumerate(fix_plans):
                is_bad = "别这样写" in plan.text
                x_pos = float(MARGIN_L.inches) if is_bad else MARGIN_L.inches + col_w + 0.25
                _draw_fix_card_plan(
                    x_pos,
                    y,
                    col_w,
                    plan,
                    warning=is_bad,
                    bold=not is_bad,
                )
        elif fix_plans:
            plan = fix_plans[0]
            is_bad = "别这样写" in plan.text
            _draw_fix_card_plan(
                float(MARGIN_L.inches),
                y,
                float(CONTENT_W.inches),
                plan,
                warning=is_bad,
                bold=not is_bad,
            )

    def vocab_table_slide(
        self,
        title: str,
        tier: str,
        rows: list[dict],
        columns: list[str],
        *,
        badge: str | None = None,
    ) -> None:
        slide = self._blank()
        top = self._header(slide, title, section="词汇")

        header_map = {"english": "词块", "chinese": "释义", "example": "例句"}
        headers = [header_map[c] for c in columns]
        n_cols = len(columns)
        col_fracs = [0.38, 0.62] if n_cols == 2 else [0.26, 0.20, 0.54]
        content_top = top + 0.06
        budget = LAYOUT_REGISTRY["vocab_table"]
        row_values: list[list[str]] = [[row.get(c, "") for c in columns] for row in rows]
        row_heights = solve_table_row_heights(row_values, col_fracs, header=True)

        tbl_shape = slide.shapes.add_table(
            len(row_heights), n_cols, MARGIN_L, Inches(content_top), CONTENT_W, Inches(sum(row_heights))
        )
        tbl = tbl_shape.table
        for i, frac in enumerate(col_fracs):
            tbl.columns[i].width = Inches(CONTENT_W.inches * frac)
        for ri, h in enumerate(row_heights):
            tbl.rows[ri].height = Inches(h)
        self._style_table_header(tbl, headers)
        for ri, vals in enumerate(row_values, start=1):
            rh = row_heights[ri]
            for c, val in enumerate(vals):
                col_w = float(CONTENT_W.inches) * col_fracs[c]
                kind: str = "primary" if columns[c] in ("english", "example") else "secondary"
                max_pt = budget.max_primary_pt if kind == "primary" else budget.max_secondary_pt
                cell_fit = fit_typography(
                    val,
                    max(1.5, col_w - 0.12),
                    max(0.35, rh - 0.20),
                    min_pt=budget.min_pt,
                    max_pt=max_pt,
                )
                pt = cell_fit.font_pt
                sp = cell_fit.line_spacing
                cell = tbl.cell(ri, c)
                cell.text = val
                _configure_text_frame(cell.text_frame, anchor=MSO_ANCHOR.TOP)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(pt)
                    p.line_spacing = sp
                    if columns[c] == "english":
                        p.font.name = FONT_ESSAY
                        p.font.bold = True
                        p.font.color.rgb = INK
                    elif columns[c] == "example":
                        p.font.name = FONT_ESSAY
                        p.font.color.rgb = INK
                    else:
                        p.font.name = FONT_UI
                        p.font.color.rgb = MUTED
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = G.ZONE_SURFACE
        tbl_shape.height = Inches(sum(row_heights))


def render_v2_deck(slides: list[dict], output: Path) -> Path:
    slides = pack_slides(slides)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    builder = SlideBuilderV2(prs)
    for spec in slides:
        kind = spec["type"]
        builder._set_spec_context(spec)
        solved = solve_layout(spec)
        if kind == "title":
            builder.title_slide(
                spec["title"],
                spec.get("subtitle", ""),
                spec.get("body") or [],
                poster_lines=spec.get("poster_lines"),
                task_tag=spec.get("task_tag", ""),
            )
        elif kind == "title_poster":
            builder.content_slide(
                spec.get("title", "海报示意"),
                spec.get("poster_lines") or [],
                panel=True,
            )
        elif kind == "roadmap":
            pass
        elif kind == "divider":
            continue
        elif kind == "essay":
            if solved.essay_split_parts:
                for idx, part in enumerate(solved.essay_split_parts):
                    part_title = (
                        f"{spec['title']}（{idx + 1}/{len(solved.essay_split_parts)}）"
                        if len(solved.essay_split_parts) > 1
                        else spec["title"]
                    )
                    builder.essay_slide(
                        part_title,
                        "",
                        "",
                        badge=spec.get("badge") if idx == 0 else None,
                        layout_plan=part,
                    )
            elif solved.essay:
                builder.essay_slide(
                    spec["title"],
                    spec.get("essay_text", ""),
                    spec.get("annotation", ""),
                    badge=spec.get("badge"),
                    layout_plan=solved.essay,
                )
            else:
                builder.essay_slide(
                    spec["title"],
                    spec.get("essay_text", ""),
                    spec.get("annotation", ""),
                    badge=spec.get("badge"),
                )
        elif kind == "table":
            builder.table_slide(spec["title"], spec["headers"], spec["rows"])
        elif kind == "peel":
            builder.peel_slide(
                spec["title"],
                spec["points"],
                layout=spec.get("layout", "dual"),
            )
        elif kind == "phrase_table":
            builder.phrase_table_slide(
                spec["title"],
                spec["table"],
                badge=spec.get("badge"),
                part=spec.get("part", "full"),
            )
        elif kind == "vocab_table":
            builder.vocab_table_slide(
                spec["title"],
                spec.get("tier", ""),
                spec["rows"],
                spec["columns"],
                badge=spec.get("badge"),
            )
        else:
            builder.content_slide(
                spec["title"],
                spec["bullets"],
                badge=spec.get("badge"),
                panel=spec.get("panel", False),
                warn_panel=spec.get("warn_panel", False),
                callout=spec.get("callout"),
            )
    output.parent.mkdir(parents=True, exist_ok=True)
    _enforce_word_wrap(prs)
    _clean_pptx_dir(output)
    prs.save(output)
    return output


def build_v2_deck_with_stage3(
    stage3_path: Path,
    deck_plan_path: Path | None = None,
    *,
    vocab_max_rows: int = 6,
) -> list[dict]:
    base = build_deck_with_stage3(
        stage3_path, deck_plan_path, vocab_max_rows=vocab_max_rows
    )
    return _inject_v2_structure(base)


def render_v2_classroom(
    slides: list[dict],
    output: Path,
) -> Path:
    return render_v2_deck(slides, output)


def main(argv: list[str] | None = None) -> int:
    import argparse

    parser = argparse.ArgumentParser(description="Generate 应用文 classroom PPTX V2")
    parser.add_argument("-o", "--output", type=Path, default=DEFAULT_V2_OUT)
    parser.add_argument("--stage3", type=Path, default=None)
    parser.add_argument("--deck-plan", type=Path, default=None)
    args = parser.parse_args(argv)

    if args.stage3 and args.stage3.is_file():
        slides = build_v2_deck_with_stage3(
            args.stage3.expanduser().resolve(),
            args.deck_plan.expanduser().resolve() if args.deck_plan else None,
        )
    else:
        slides = _inject_v2_structure(expand_slide_specs(build_mental_health_deck()))
    path = render_v2_deck(slides, args.output)
    prs = Presentation(path)
    seen, violations = _collect_font_sizes(prs)
    layout = verify_deck_layout(prs)
    page_nums = _find_page_numbers(prs)

    print(f"V2 Saved: {path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Fonts: {sorted(seen)} min={min(seen) if seen else '?'}")
    print(
        f"Layout: {'pass' if layout['ok'] else 'FAIL'} "
        f"(pass1={len(layout['pass1_issues'])}, WPS_RISK={layout['wps_risk_count']}, "
        f"teach_ready={layout.get('is_teach_ready')})"
    )
    if layout.get("wps_report"):
        r = layout["wps_report"]
        print(
            f"WPS governance: critical={len(r.critical_issues)} warning={len(r.warning_issues)} "
            f"cosmetic={len(r.cosmetic_issues)} risk_score={r.risk_score:.3f}"
        )
    if layout["pass2_issues"]:
        print(f"WPS issues sample: {layout['pass2_issues'][:3]}")
    print(f"Page nums: {'none' if not page_nums else page_nums}")
    ok = not violations and layout["ok"] and not page_nums
    return 0 if ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
