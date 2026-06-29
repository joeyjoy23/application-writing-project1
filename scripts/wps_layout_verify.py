"""WPS layout verification with teaching-aware governance tiers.

Legacy verify_text_fit remains for debug; WPS governance is the release gate.
"""

from __future__ import annotations

import re
from dataclasses import dataclass, field
from typing import Any, Literal

from pptx import Presentation

from scripts.generate_classroom_pptx import verify_text_fit
from scripts.ppt_layout_fit import TABLE_MAX_ROWS_PER_PAGE, _cjk_ratio, line_count

WPS_LAYOUT_SAFETY: dict[str, float] = {
    "line_height_factor": 1.42,
    "paragraph_spacing": 1.15,
    "bold_extra_factor": 1.08,
    "chinese_extra_factor": 1.05,
    "table_row_padding": 1.20,
}

TABLE_ROW_HEIGHT_MULT = 1.25
BULLET_EXTRA_FACTOR = 1.06
TABLE_CELL_EXTRA_FACTOR = 1.10
MIN_MARGIN_RATIO = 0.12
OVERFLOW_TOLERANCE = 1.02
WARNING_RATIO_PER_SLIDE = 0.20

SECTION_LABELS = frozenset(
    {"审题", "句型", "范文", "思维", "路线", "导入", "批改", "迁移", "词汇", "PEEL", "句型库"}
)
UI_LABELS = frozenset({"抓重点", "本题"})

IssueLevel = Literal["critical", "warning", "cosmetic"]

_CJK_CHAR_RE = re.compile(r"[\u4e00-\u9fff，。；：！？、「」]")
_WORD_RE = re.compile(r"\S+|\s+")


def safe_padding_inches(font_pt: float) -> float:
    """Per-edge vertical padding: max(14px, font_size * 0.35) in inches."""
    font_in = max(float(font_pt), 1.0) / 72.0
    return max(14.0 / 72.0, font_in * 0.35)


def _wps_chars_per_line(text: str, width_inches: float, font_pt: float) -> float:
    ratio = _cjk_ratio(text)
    base = 26.0 * ratio + 50.0 * (1.0 - ratio)
    scale = (max(0.5, width_inches) * 72.0) / (max(font_pt, 8.0) * 1.08)
    return max(6.0, base * scale / 6.8)


def _wps_wrap_line_count(text: str, width_inches: float, font_pt: float, *, bullet: bool) -> float:
    clean = (text or "").strip()
    if not clean:
        return 0.35

    cpl = _wps_chars_per_line(clean, width_inches, font_pt)
    total = 0.0
    for para in clean.split("\n"):
        para = para.strip()
        if not para:
            total += 0.35
            continue

        if bullet or para.startswith(("•", "①", "②", "③", "④", "★", "→")):
            body = para.lstrip("•①②③④★→ ").strip()
            if not body:
                total += 0.35
                continue
            para = body

        ratio = _cjk_ratio(para)
        if ratio > 0.55:
            total += max(1.0, len(para) / cpl)
            continue

        line_len = 0.0
        line_count_local = 1.0
        for token in _WORD_RE.findall(para):
            if token.isspace():
                continue
            tok_len = sum(1.2 if _CJK_CHAR_RE.search(ch) else 1.0 for ch in token)
            if line_len > 0 and line_len + tok_len > cpl:
                line_count_local += 1.0
                line_len = tok_len
            else:
                line_len += tok_len
        total += max(1.0, line_count_local)

    return total


def estimate_wps_text_height(
    text: str,
    font_size: float,
    style: dict[str, bool] | None = None,
    *,
    width_inches: float = 5.0,
) -> float:
    """Conservative WPS text block height (inches)."""
    style = style or {}
    clean = (text or "").strip()
    if not clean:
        return 0.35 / 72.0

    font_pt = max(float(font_size), 8.0)
    lh = WPS_LAYOUT_SAFETY["line_height_factor"]
    base_line_in = (font_pt / 72.0) * lh

    chinese = style.get("chinese")
    if chinese is None:
        chinese = _cjk_ratio(clean) > 0.25
    bullet = bool(style.get("bullet"))
    if not bullet and clean.startswith(("•", "①", "②", "③", "④", "★")):
        bullet = True

    mult = 1.0
    if chinese:
        mult *= WPS_LAYOUT_SAFETY["chinese_extra_factor"]
    if style.get("bold"):
        mult *= WPS_LAYOUT_SAFETY["bold_extra_factor"]
    if bullet:
        mult *= BULLET_EXTRA_FACTOR
    if style.get("table_cell"):
        mult *= TABLE_CELL_EXTRA_FACTOR

    lines = _wps_wrap_line_count(clean, width_inches, font_pt, bullet=bullet)
    para_n = max(1, len([p for p in clean.split("\n") if p.strip()]))
    para_gap = (para_n - 1) * (font_pt / 72.0) * (WPS_LAYOUT_SAFETY["paragraph_spacing"] - 1.0)
    return lines * base_line_in * mult + max(0.0, para_gap)


WPS_SAFE_FACTOR = 1.30
CARD_CONTAINER_PAD_V = 0.32
CARD_PAD_OVERFLOW_HEADROOM = 0.08


def compute_card_height(
    text: str,
    font_size: float = 28,
    width_inches: float = 11.0,
    *,
    bold: bool = False,
    bullet: bool = False,
    table_cell: bool = False,
    min_h: float = 0.56,
) -> float:
    """Bidirectional shrink-wrap: preferred = text + min pad; WPS_SAFE_FACTOR only caps overflow."""
    clean = (text or "").strip()
    if not clean:
        return min_h
    st = {
        "chinese": _cjk_ratio(clean) > 0.25,
        "bold": bold,
        "bullet": bullet or clean.startswith(("•", "①", "②", "③", "④", "★", "→")),
        "table_cell": table_cell,
    }
    text_h = estimate_wps_text_height(clean, font_size, st, width_inches=width_inches)
    pad_min = safe_padding_inches(font_size) * 2 + CARD_CONTAINER_PAD_V
    preferred = text_h + pad_min
    safe_min = text_h + pad_min
    safe_max = text_h * WPS_SAFE_FACTOR + pad_min + CARD_PAD_OVERFLOW_HEADROOM
    clamped = max(safe_min, min(preferred, safe_max))
    return max(min_h, clamped)


def compute_safe_height(
    text: str,
    font_size: float = 28,
    width_inches: float = 11.0,
    *,
    bold: bool = False,
    bullet: bool = False,
    table_cell: bool = False,
    min_h: float = 0.56,
) -> float:
    """Alias for compute_card_height (renderer + CI gate)."""
    return compute_card_height(
        text,
        font_size,
        width_inches,
        bold=bold,
        bullet=bullet,
        table_cell=table_cell,
        min_h=min_h,
    )


def verify_heuristic_text_height(
    text: str,
    font_size: float,
    *,
    width_inches: float = 5.0,
    line_spacing: float = 1.12,
) -> float:
    return line_count(text, width_inches, font_size) * (font_size / 72.0) * line_spacing


def _infer_style_from_text(text: str, *, table_cell: bool = False, bold_hint: bool = False) -> dict[str, bool]:
    return {
        "chinese": _cjk_ratio(text) > 0.25,
        "bold": bold_hint,
        "bullet": text.strip().startswith(("•", "①", "②", "③", "④", "★", "→")),
        "table_cell": table_cell,
    }


def _paragraph_font_pt(tf) -> list[tuple[str, float, bool]]:
    rows: list[tuple[str, float, bool]] = []
    for para in tf.paragraphs:
        t = para.text
        if not t.strip():
            continue
        pt = 26.0
        bold = False
        if para.font.size:
            pt = float(para.font.size.pt)
        if para.font.bold:
            bold = True
        for run in para.runs:
            if run.font.size:
                pt = max(pt, float(run.font.size.pt))
            if run.font.bold:
                bold = True
        rows.append((t, pt, bold))
    return rows


@dataclass
class WpsElementDebug:
    slide_id: int
    element_id: str
    estimated_height: float
    shape_height: float
    wps_safe_height: float
    margin_ratio: float
    kind: str = "text_frame"
    level: IssueLevel = "warning"
    cosmetic_reason: str = ""


@dataclass
class WpsSlideDebug:
    slide_id: int
    elements: list[WpsElementDebug] = field(default_factory=list)
    info: list[str] = field(default_factory=list)


@dataclass
class WPSLayoutIssue:
    level: IssueLevel
    slide_id: int
    element_id: str
    message: str
    estimated_height: float = 0.0
    shape_height: float = 0.0
    margin_ratio: float = 0.0
    reason: str = ""


@dataclass
class WPSLayoutReport:
    critical_issues: list[WPSLayoutIssue] = field(default_factory=list)
    warning_issues: list[WPSLayoutIssue] = field(default_factory=list)
    cosmetic_issues: list[WPSLayoutIssue] = field(default_factory=list)
    is_teach_ready: bool = True
    risk_score: float = 0.0
    total_slides: int = 0

    @property
    def wps_risk_overflow_count(self) -> int:
        """Critical-only count (CI gate)."""
        return len(self.critical_issues)


def _wps_safe_envelope(text_height: float, font_pt: float) -> float:
    pad = safe_padding_inches(font_pt)
    return text_height + 2.0 * pad


def _margin_ratio(shape_height: float, text_height: float) -> float:
    if text_height <= 1e-6:
        return 1.0
    return (shape_height - text_height) / text_height


def _shape_metrics(shape) -> tuple[float, float, float, float]:
    top = shape.top.inches if shape.top else 0.0
    left = shape.left.inches if shape.left else 0.0
    height = shape.height.inches if shape.height else 0.0
    width = shape.width.inches if shape.width else 0.0
    return top, left, height, width


def classify_wps_element(
    *,
    text: str,
    shape_top: float,
    shape_left: float,
    shape_h: float,
    shape_w: float,
    est: float,
    safe_h: float,
    margin_ratio: float,
    kind: str,
) -> tuple[IssueLevel, str]:
    """Classify one text-bearing element into critical / warning / cosmetic."""
    clean = (text or "").strip()
    if not clean:
        return "cosmetic", "empty_text"

    if shape_h < 0.12 or clean in {"↓", "→", "↑", "▼", "▶"}:
        return "cosmetic", "decorative_or_flow_arrow"

    if clean in {"① 先选", "② 再讲理由", "① 选择", "② 理由"}:
        return "cosmetic", "peel_card_label"

    if re.match(r"^[①②③④]", clean) and len(clean) <= 12 and shape_h <= 0.62:
        return "cosmetic", "peel_card_label"

    if shape_top < 0.25 and shape_left < 0.65 and shape_w <= 2.2 and shape_h <= 0.58:
        if clean in SECTION_LABELS or len(clean) <= 5:
            return "cosmetic", "section_pill"

    if clean in UI_LABELS and shape_h <= 0.58 and shape_w <= 2.2:
        return "cosmetic", "ui_label"

    if clean.startswith("本课任务"):
        return "cosmetic", "task_pill"

    if shape_top < 0.50 and shape_w >= 5.0 and len(clean) < 55:
        return "cosmetic", "cover_title"

    if 0.72 <= shape_top <= 0.88 and shape_h <= 0.82 and shape_w >= 9.5 and len(clean) < 90:
        return "cosmetic", "slide_title"

    if shape_h <= 0.58 and shape_w <= 4.5 and len(clean) < 22:
        if clean.startswith("Stage ") or ("·" in clean and not clean.startswith("•")):
            return "cosmetic", "badge_pill"
        if not clean.startswith("•") and "。" not in clean and "；" not in clean:
            if len(clean) <= 14 and not any(k in clean for k in ("假如", "Poster", "Dear")):
                return "cosmetic", "short_pill"

    actual_overflow = est > shape_h * OVERFLOW_TOLERANCE
    overflow_severity = (est - shape_h) / max(shape_h, 0.08) if actual_overflow else 0.0
    tight_padding = est <= shape_h and margin_ratio < MIN_MARGIN_RATIO

    is_bullet = clean.startswith(("•", "①", "★", "→")) or "Poster" in clean
    required_h = compute_safe_height(
        clean,
        26.0,
        shape_w,
        bold="改成" in clean or clean.startswith("★"),
        bullet=is_bullet,
        table_cell=kind == "table_cell",
        min_h=0.12 if kind == "table_cell" else 0.56,
    )
    under_allocated = shape_h < required_h * 0.96

    if kind == "table_cell" and shape_h < 0.28:
        return "critical", "table_row_collapsed"

    if under_allocated and actual_overflow:
        if kind == "table_cell":
            return "critical", "table_cell_under_allocated"
        if "Dear " in clean or (len(clean) > 220 and shape_h > 0.8):
            return "critical", "essay_truncation_risk"
        return "critical", "text_overflow_shape"

    if actual_overflow:
        if kind == "table_cell":
            return "warning", "table_cell_tight"
        if is_bullet:
            return "warning", "content_card_tight"
        return "warning", "wps_model_overflow"

    if tight_padding or (est > safe_h and est <= shape_h):
        if kind == "table_cell":
            return "warning", "table_cell_tight"
        if is_bullet:
            return "warning", "content_card_tight"
        return "warning", "low_margin_ratio"

    return "cosmetic", "within_safe_envelope"


def _issue_from_element(
    slide_id: int,
    element_id: str,
    level: IssueLevel,
    reason: str,
    *,
    est: float,
    shape_h: float,
    safe_h: float,
    margin_ratio: float,
    preview: str,
) -> WPSLayoutIssue:
    tag = "WPS_CRITICAL" if level == "critical" else "WPS_WARNING"
    return WPSLayoutIssue(
        level=level,
        slide_id=slide_id,
        element_id=element_id,
        message=(
            f"{tag}: slide {slide_id} {element_id} ({reason}) "
            f"est={est:.2f}\" shape={shape_h:.2f}\" safe={safe_h:.2f}\" "
            f"margin_ratio={margin_ratio:.2f} [{preview[:40]}...]"
        ),
        estimated_height=est,
        shape_height=shape_h,
        margin_ratio=margin_ratio,
        reason=reason,
    )


def compute_risk_score(report: WPSLayoutReport) -> float:
    slides = max(report.total_slides, 1)
    crit = len(report.critical_issues)
    warn = len(report.warning_issues)
    raw = (crit * 1.0 + warn * 0.25) / slides
    return min(1.0, round(raw, 3))


def is_teach_ready(report: WPSLayoutReport) -> bool:
    if report.critical_issues:
        return False
    slides = max(report.total_slides, 1)
    return len(report.warning_issues) / slides < WARNING_RATIO_PER_SLIDE


def verify_wps_layout_safety(prs_or_slides: Presentation) -> tuple[WPSLayoutReport, list[WpsSlideDebug]]:
    """Scan deck; classify each text element into critical / warning / cosmetic."""
    prs = prs_or_slides
    report = WPSLayoutReport(total_slides=len(prs.slides))
    slide_debug: list[WpsSlideDebug] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        dbg = WpsSlideDebug(slide_id=slide_idx)
        table_row_counts: list[int] = []

        for shape_idx, shape in enumerate(slide.shapes, start=1):
            if shape.has_table:
                table = shape.table
                table_row_counts.append(len(table.rows))
                for ri, row in enumerate(table.rows):
                    row_h = row.height.inches if row.height else shape.height.inches / max(len(table.rows), 1)
                    eff_row_h = row_h * TABLE_ROW_HEIGHT_MULT
                    for ci, cell in enumerate(row.cells):
                        tf = cell.text_frame
                        if not any(p.text.strip() for p in tf.paragraphs):
                            continue
                        col_w = (
                            table.columns[ci].width.inches
                            if table.columns[ci].width
                            else shape.width.inches / max(len(table.columns), 1)
                        )
                        paras = _paragraph_font_pt(tf)
                        est = 0.0
                        max_pt = 26.0
                        for text, pt, bold in paras:
                            max_pt = max(max_pt, pt)
                            st = _infer_style_from_text(text, table_cell=True, bold_hint=bold)
                            est += estimate_wps_text_height(text, pt, st, width_inches=col_w)
                        safe_h = _wps_safe_envelope(est, max_pt)
                        mr = _margin_ratio(eff_row_h, est)
                        elem_id = f"table_r{ri + 1}c{ci + 1}"
                        full_text = tf.text
                        top, left, _, _ = _shape_metrics(shape)
                        level, reason = classify_wps_element(
                            text=full_text,
                            shape_top=top,
                            shape_left=left,
                            shape_h=eff_row_h,
                            shape_w=col_w,
                            est=est,
                            safe_h=safe_h,
                            margin_ratio=mr,
                            kind="table_cell",
                        )
                        dbg.elements.append(
                            WpsElementDebug(
                                slide_id=slide_idx,
                                element_id=elem_id,
                                estimated_height=est,
                                shape_height=eff_row_h,
                                wps_safe_height=safe_h,
                                margin_ratio=mr,
                                kind="table_cell",
                                level=level,
                                cosmetic_reason=reason if level == "cosmetic" else "",
                            )
                        )
                        if level == "cosmetic":
                            report.cosmetic_issues.append(
                                _issue_from_element(
                                    slide_idx,
                                    elem_id,
                                    level,
                                    reason,
                                    est=est,
                                    shape_h=eff_row_h,
                                    safe_h=safe_h,
                                    margin_ratio=mr,
                                    preview=full_text,
                                )
                            )
                            continue
                        issue = _issue_from_element(
                            slide_idx,
                            elem_id,
                            level,
                            reason,
                            est=est,
                            shape_h=eff_row_h,
                            safe_h=safe_h,
                            margin_ratio=mr,
                            preview=full_text,
                        )
                        if level == "critical":
                            report.critical_issues.append(issue)
                        else:
                            report.warning_issues.append(issue)
                continue

            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if not any(p.text.strip() for p in tf.paragraphs):
                continue

            top, left, shape_h, shape_w = _shape_metrics(shape)
            paras = _paragraph_font_pt(tf)
            est = 0.0
            max_pt = 26.0
            for text, pt, bold in paras:
                max_pt = max(max_pt, pt)
                st = _infer_style_from_text(text, bold_hint=bold)
                est += estimate_wps_text_height(text, pt, st, width_inches=shape_w)
            safe_h = _wps_safe_envelope(est, max_pt)
            mr = _margin_ratio(shape_h, est)
            elem_id = f"shape_{shape_idx}"
            full_text = tf.text
            level, reason = classify_wps_element(
                text=full_text,
                shape_top=top,
                shape_left=left,
                shape_h=shape_h,
                shape_w=shape_w,
                est=est,
                safe_h=safe_h,
                margin_ratio=mr,
                kind="text_frame",
            )
            dbg.elements.append(
                WpsElementDebug(
                    slide_id=slide_idx,
                    element_id=elem_id,
                    estimated_height=est,
                    shape_height=shape_h,
                    wps_safe_height=safe_h,
                    margin_ratio=mr,
                    kind="text_frame",
                    level=level,
                    cosmetic_reason=reason if level == "cosmetic" else "",
                )
            )
            if level == "cosmetic":
                report.cosmetic_issues.append(
                    _issue_from_element(
                        slide_idx,
                        elem_id,
                        level,
                        reason,
                        est=est,
                        shape_h=shape_h,
                        safe_h=safe_h,
                        margin_ratio=mr,
                        preview=full_text,
                    )
                )
                continue
            issue = _issue_from_element(
                slide_idx,
                elem_id,
                level,
                reason,
                est=est,
                shape_h=shape_h,
                safe_h=safe_h,
                margin_ratio=mr,
                preview=full_text,
            )
            if level == "critical":
                report.critical_issues.append(issue)
            else:
                report.warning_issues.append(issue)

        if table_row_counts:
            max_rows = max(table_row_counts)
            if max_rows > TABLE_MAX_ROWS_PER_PAGE:
                dbg.info.append(
                    f"table has {max_rows} rows (>{TABLE_MAX_ROWS_PER_PAGE}, informational)"
                )

        slide_debug.append(dbg)

    report.risk_score = compute_risk_score(report)
    report.is_teach_ready = is_teach_ready(report)
    return report, slide_debug


def verify_deck_layout(prs: Presentation) -> dict[str, Any]:
    """Dual validation: legacy verify_text_fit + WPS governance layer."""
    pass1 = verify_text_fit(prs)
    wps_report, slide_debug = verify_wps_layout_safety(prs)
    critical = wps_report.critical_issues
    ok = not pass1 and not critical
    pass2_messages = [i.message for i in critical]
    warning_messages = [i.message for i in wps_report.warning_issues]
    return {
        "ok": ok,
        "pass1_issues": pass1,
        "pass2_issues": pass2_messages,
        "warning_issues": warning_messages,
        "wps_report": wps_report,
        "wps_risk_count": len(critical),
        "is_teach_ready": wps_report.is_teach_ready,
        "risk_score": wps_report.risk_score,
        "slide_debug": slide_debug,
    }


def print_wps_layout_report(prs: Presentation) -> None:
    """Print governance summary and per-element diagnostics."""
    report, slide_debug = verify_wps_layout_safety(prs)
    print("WPS Layout Governance Summary")
    print(f"  critical:          {len(report.critical_issues)}")
    print(f"  warning:           {len(report.warning_issues)}")
    print(f"  cosmetic (ignored): {len(report.cosmetic_issues)}")
    print(f"  teach_ready:       {report.is_teach_ready}")
    print(f"  risk_score:        {report.risk_score:.3f}")
    print(f"  WPS_RISK_OVERFLOW: {report.wps_risk_overflow_count} (critical only)")
    for issue in report.critical_issues[:8]:
        print(f"  [CRITICAL] {issue.message}")
    for issue in report.warning_issues[:5]:
        print(f"  [WARNING]  {issue.message}")
    if len(report.warning_issues) > 5:
        print(f"  ... +{len(report.warning_issues) - 5} warnings")
    for sd in slide_debug:
        for info in sd.info:
            print(f"  [info] slide {sd.slide_id}: {info}")
        for el in sd.elements:
            if el.level == "cosmetic":
                continue
            print(
                f"slide {el.slide_id} {el.element_id} [{el.level}]: "
                f"estimated_height={el.estimated_height:.3f}\" "
                f"shape_height={el.shape_height:.3f}\" "
                f"wps_safe_height={el.wps_safe_height:.3f}\" "
                f"margin_ratio={el.margin_ratio:.3f}"
            )


def print_wps_layout_debug(prs: Presentation) -> None:
    """Alias for governance report output."""
    print_wps_layout_report(prs)
