#!/usr/bin/env python
"""Full PPT visual/structure audit — report only."""
from __future__ import annotations

import json
import sys
from collections import defaultdict
from pathlib import Path

_PROJECT = Path(__file__).resolve().parents[1]
if str(_PROJECT) not in sys.path:
    sys.path.insert(0, str(_PROJECT))

from pptx import Presentation

PPTX = Path(r"D:\Downloads\ppt-work\mental_health_classroom.pptx")
WORK = Path(r"D:\Downloads\ppt-work")
OUT = Path(r"D:\Downloads\ppt-work\audit-full-report.txt")


def _rgb_hex(shape) -> str:
    try:
        if shape.fill.type is not None:
            c = shape.fill.fore_color.rgb
            if c:
                return f"#{c[0]:02X}{c[1]:02X}{c[2]:02X}"
    except Exception:
        pass
    return ""


def _text_shapes(slide) -> list[dict]:
    out: list[dict] = []
    for i, sh in enumerate(slide.shapes, 1):
        if not sh.has_text_frame:
            continue
        t = (sh.text_frame.text or "").strip()
        if not t:
            continue
        out.append(
            {
                "idx": i,
                "top": sh.top.inches if sh.top else 0,
                "h": sh.height.inches if sh.height else 0,
                "w": sh.width.inches if sh.width else 0,
                "left": sh.left.inches if sh.left else 0,
                "text": t.replace("\n", " | "),
                "fill": _rgb_hex(sh),
            }
        )
    return out


def main() -> None:
    lines: list[str] = []
    if not PPTX.is_file():
        lines.append("ERROR: PPTX not found")
        OUT.write_text("\n".join(lines), encoding="utf-8")
        print(lines[0])
        return

    prs = Presentation(str(PPTX))
    issues: dict[str, list] = defaultdict(list)
    slide_rows: list[tuple] = []

    for si, slide in enumerate(prs.slides, 1):
        ts = _text_shapes(slide)
        if not ts:
            slide_rows.append((si, "(empty)", 1.0, 0, []))
            issues["empty_slide"].append(si)
            continue
        bottom = max(x["top"] + x["h"] for x in ts)
        blank = max(0, (7.5 - bottom) / 7.5)
        title = ts[0]["text"][:70]
        colors = {_rgb_hex(sh) for sh in slide.shapes if _rgb_hex(sh)}
        slide_rows.append((si, title, blank, len(colors), ts))

        if blank > 0.30:
            issues["large_blank"].append((si, blank, title))

        pills = [x for x in ts if x["top"] < 0.25 and x["h"] <= 0.55 and x["w"] <= 2.5]
        headers = [x for x in ts if 0.72 <= x["top"] <= 0.88 and x["w"] >= 9]
        for p in pills:
            for h in headers:
                pt, ht = p["text"], h["text"]
                if pt in ht or ht.startswith(pt.split(" ·")[0]):
                    issues["title_redundant"].append((si, pt, ht[:55]))

        for x in ts:
            if x["text"].startswith("•") and x["h"] > 0.85 and len(x["text"]) < 120:
                issues["tall_bullet_card"].append((si, x["h"], x["text"][:60]))

        for x in ts:
            if x["h"] < 0.4 or len(x["text"]) < 10:
                continue
            for sh in slide.shapes:
                if sh.shape_type != 1 or sh.has_text_frame:
                    continue
                rt = sh.top.inches if sh.top else 0
                rh = sh.height.inches if sh.height else 0
                if abs(x["top"] - rt) < 0.25 and rh > x["h"] * 1.45 and rh > 0.65:
                    issues["card_text_mismatch"].append(
                        (si, rh / x["h"], x["text"][:50], rh, x["h"])
                    )

        if len(colors) > 7:
            issues["many_fill_colors"].append((si, len(colors), title))

    lines.append("=" * 70)
    lines.append("PPT 全量审计报告（只读，未改代码）")
    lines.append(f"文件: {PPTX}")
    lines.append(f"页数: {len(prs.slides)}")
    lines.append("")

    lines.append("--- 逐页概览 ---")
    for si, title, blank, ncolors, ts in slide_rows:
        flag = ""
        if blank > 0.30:
            flag += " [空白>30%]"
        if any(si == r[0] for r in issues["title_redundant"]):
            flag += " [标题重复]"
        lines.append(f"Slide {si:2d} blank={blank:.0%} colors={ncolors} | {title}{flag}")

    lines.append("")
    lines.append("--- 问题统计 ---")
    for k, v in sorted(issues.items(), key=lambda x: -len(x[1])):
        lines.append(f"  {k}: {len(v)}")

    def _section(name: str, key: str, fmt):
        if not issues[key]:
            return
        lines.append("")
        lines.append(f"--- {name} ---")
        for row in issues[key][:20]:
            lines.append("  " + fmt.format(*row))
        if len(issues[key]) > 20:
            lines.append(f"  ... +{len(issues[key]) - 20} more")

    _section("P1 大片空白（底部>30%）", "large_blank", "Slide {} blank={:.0%} | {}")
    _section("P2 标题重复（pill + 主标题）", "title_redundant", "Slide {} pill='{}' header='{}'")
    _section("P3 卡片/框过高（单行 bullet）", "tall_bullet_card", "Slide {} h={:.2f}\" | {}")
    _section("P4 框与文字高度不匹配", "card_text_mismatch", "Slide {} ratio={:.1f}x | {} (box {:.2f} text {:.2f})")

    # spec pipeline
    lines.append("")
    lines.append("--- 编译/打包链 ---")
    try:
        from scripts.one_click_classroom_ppt import _build_v2_slides
        from scripts.generate_classroom_pptx_v2 import _inject_v2_structure
        from scripts.ppt_layout_fit import (
            expand_content_slides,
            expand_essay_slides,
            expand_peel_slides,
            expand_title_slides,
            pack_slides,
        )

        stage3 = WORK / "stage3.json"
        export = WORK / "export-data.json"
        if stage3.is_file():
            slides = _build_v2_slides(
                stage3,
                None,
                vocab_max_rows=6,
                use_custom_plan=False,
                export_data_path=export if export.is_file() else None,
                preset="70min",
                script_path=None,
                template_id="dual_poster_opinion",
            )
            slides = _inject_v2_structure(slides)
            expanded = expand_content_slides(
                expand_peel_slides(expand_essay_slides(expand_title_slides(slides)))
            )
            packed = pack_slides(expanded)
            lines.append(f"  compiled={len(slides)} expanded={len(expanded)} packed={len(packed)} pptx={len(prs.slides)}")
            lines.append("  前 8 页 spec:")
            for i, s in enumerate(packed[:8], 1):
                lines.append(f"    {i:2d} {s.get('type', '?'):14s} {(s.get('title') or '')[:48]}")
            if packed:
                t0 = packed[0]
                t1 = packed[1] if len(packed) > 1 else {}
                if t0.get("type") == "title" and t1.get("title", "").startswith("导入"):
                    lines.append("  [合并候选] Slide1 封面 + Slide2 导入·真题展示 → 可合并为一页")
                if any(s.get("type") == "title_poster" for s in packed[:3]):
                    lines.append("  [拆分残留] expand_title_slides 仍产出 title_poster 独立页")
    except Exception as exc:
        lines.append(f"  spec chain error: {exc}")

    lines.append("")
    lines.append("--- 根因层归类（建议修复方向，本次未改） ---")
    lines.append("  A. Script/页序: cover_meta 与 question_stem 分两页；pitfall_contrast 标题含「审题」")
    lines.append("  B. expand_title_slides: poster/stem 拆成 title_poster 增加空白页")
    lines.append("  C. Renderer header: _section_tag + _header 双行标题占 ~1.6\" 垂直空间")
    lines.append("  D. Layout: compute_safe_height 偏保守 → 单行卡片过高；禁止拆页后内容稀疏")
    lines.append("  E. 配色: 若仍见紫粉，请 Reboot Streamlit / 本地重生成并核对 UI_BUILD_TAG")
    lines.append("")
    lines.append("--- WPS 校验（参考） ---")
    try:
        from scripts.wps_layout_verify import verify_deck_layout

        layout = verify_deck_layout(prs)
        r = layout.get("wps_report")
        lines.append(
            f"  pass1={len(layout['pass1_issues'])} critical={layout['wps_risk_count']} "
            f"warning={len(r.warning_issues) if r else '?'} teach_ready={layout.get('is_teach_ready')}"
        )
    except Exception as exc:
        lines.append(f"  wps check error: {exc}")

    lines.append("")
    lines.append("END")
    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {OUT} ({len(lines)} lines)")


if __name__ == "__main__":
    main()
