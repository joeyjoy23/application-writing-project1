#!/usr/bin/env python
"""Generate editable classroom PPTX with integrated Stage4 pedagogy."""

from __future__ import annotations

import argparse
import re
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

_SCRIPTS_DIR = Path(__file__).resolve().parent
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

from essay_format import classroom_body_paragraphs, essay_layout_for_length
from deck_plan import load_deck_plan, replace_stage3_in_deck, stage3_specs_from_plan
from pptx_click_reveal import apply_click_reveal
from scripts.ppt_layout_fit import effective_text_area, line_count

PRIMARY = RGBColor(0x63, 0x66, 0xF1)
ACCENT = RGBColor(0x14, 0xB8, 0xA6)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
BODY = RGBColor(0x1E, 0x29, 0x3B)
SECONDARY = RGBColor(0x64, 0x74, 0x8B)
WARNING = RGBColor(0xEF, 0x44, 0x44)
PANEL_BG = RGBColor(0xF5, 0xF3, 0xFF)
PANEL_SOFT = RGBColor(0xF8, 0xFA, 0xFC)
PANEL_ALT = RGBColor(0xEC, 0xFE, 0xFF)
WARNING_PANEL = RGBColor(0xFE, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL_BORDER = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.65)
CONTENT_W = Inches(12.0)
BOTTOM_Y = 7.32
HEADER_H = Inches(0.82)
TEXT_PAD_H = Inches(0.15)
TEXT_PAD_V = Inches(0.12)
COVER_PANEL_PAD = 0.22

FONT_UI = "Microsoft YaHei UI"
FONT_ESSAY = "Times New Roman"

# STRICT: minimum 26pt everywhere
FONT_TITLE = Pt(32)
FONT_SECTION = Pt(28)
FONT_BODY = Pt(28)
FONT_ESSAY_SIZE = Pt(26)
FONT_BADGE = Pt(26)
FONT_TABLE = Pt(26)
FONT_ANNOTATION = Pt(26)

MIN_FONT_PT = 26
MAX_BULLETS_PER_SLIDE = 5
MAX_BODY_WEIGHT = 6.0
_CHARS_PER_LINE_BODY = 32
_CHARS_PER_LINE_CJK = 28
_CHARS_PER_LINE_LATIN = 52


def _line_count(text: str, chars_per_line: int) -> int:
    if not text:
        return 1
    total = 0
    for line in text.split("\n"):
        total += max(1, (len(line) + chars_per_line - 1) // chars_per_line)
    return total


def _table_row_height(values: list[str], col_fracs: list[float]) -> float:
    """Estimate row height (inches) from wrapped cell text."""
    caps = [
        max(12, int(_CHARS_PER_LINE_CJK * frac * 0.85)) for frac in col_fracs
    ]
    lines = max(_line_count(val, cap) for val, cap in zip(values, caps, strict=True))
    return min(1.45, 0.40 + lines * 0.26)


def _configure_text_frame(
    tf,
    *,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    pad_h: int = 6,
    pad_v: int = 4,
) -> None:
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(pad_h)
    tf.margin_right = Pt(pad_h)
    tf.margin_top = Pt(pad_v)
    tf.margin_bottom = Pt(pad_v)


def _textbox_width() -> Inches:
    return CONTENT_W - TEXT_PAD_H * 2


def _textbox_left() -> Inches:
    return MARGIN_L + TEXT_PAD_H


def _paragraph_sizes_and_spacing(tf) -> tuple[list[int], list[float], list[str]]:
    texts: list[str] = []
    sizes: list[int] = []
    spacings: list[float] = []
    for para in tf.paragraphs:
        if not para.text.strip():
            continue
        texts.append(para.text)
        pt = MIN_FONT_PT
        if para.font.size:
            pt = max(MIN_FONT_PT, int(round(para.font.size.pt)))
        else:
            for run in para.runs:
                if run.font.size:
                    pt = max(pt, int(round(run.font.size.pt)))
        if pt == MIN_FONT_PT and not para.runs:
            pt = MIN_FONT_PT
        sizes.append(pt)
        spacings.append(para.line_spacing if para.line_spacing else 1.12)
    return sizes, spacings, texts


def _para_space_after_pt(tf, para_index: int) -> float:
    paras = [p for p in tf.paragraphs if p.text.strip()]
    if para_index >= len(paras) - 1:
        return 0.0
    p = paras[para_index]
    if p.space_after:
        return float(p.space_after.pt)
    return 0.0


def _frame_text_height(
    tf,
    width_inches: float,
    height_inches: float,
    *,
    pad_h_pt: float = 6,
    pad_v_pt: float = 4,
    for_verify: bool = False,
) -> tuple[float, float]:
    """Return (needed height inches, capacity height inches).

    Planning uses WPS safety factors; verify uses margins only + small tolerance.
    """
    if for_verify:
        from scripts.ppt_layout_fit import WPS_SAFETY_FACTOR

        eff_w = max(1.0, width_inches - 2 * pad_h_pt / 72.0) * WPS_SAFETY_FACTOR
        eff_h = max(0.25, height_inches - 2 * pad_v_pt / 72.0) * WPS_SAFETY_FACTOR
    else:
        eff_w, eff_h = effective_text_area(
            width_inches, height_inches, pad_h_pt=pad_h_pt, pad_v_pt=pad_v_pt
        )
    sizes, spacings, texts = _paragraph_sizes_and_spacing(tf)
    if not texts:
        return 0.0, eff_h
    needed = 0.0
    for i, (text, pt, sp) in enumerate(zip(texts, sizes, spacings, strict=True)):
        needed += line_count(text, eff_w, pt) * (pt / 72.0) * sp
        gap = _para_space_after_pt(tf, i)
        if gap:
            needed += gap / 72.0
    return needed, eff_h


def verify_text_fit(prs: Presentation) -> list[str]:
    issues: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for ri, row in enumerate(table.rows):
                    row_h = row.height.inches if row.height else shape.height.inches / max(
                        len(table.rows), 1
                    )
                    for ci, cell in enumerate(row.cells):
                        tf = cell.text_frame
                        if not any(p.text.strip() for p in tf.paragraphs):
                            continue
                        col_w = (
                            table.columns[ci].width.inches
                            if table.columns[ci].width
                            else shape.width.inches / max(len(table.columns), 1)
                        )
                        needed, cap = _frame_text_height(
                            tf, col_w, row_h, pad_h_pt=6, pad_v_pt=4, for_verify=True
                        )
                        if needed > cap * 1.02:
                            issues.append(
                                f"slide {slide_idx}: table cell overflow "
                                f"(need ~{needed:.2f}\", cap ~{cap:.2f}\")"
                            )
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            if not any(p.text.strip() for p in tf.paragraphs):
                continue
            # Skip decorative pills / section numbers (single short label).
            if shape.height.inches < 0.58 and len(tf.text.strip()) < 24:
                continue
            pad_h = tf.margin_left.pt if tf.margin_left else 6
            pad_v = tf.margin_top.pt if tf.margin_top else 4
            needed, cap = _frame_text_height(
                tf,
                shape.width.inches,
                shape.height.inches,
                pad_h_pt=pad_h,
                pad_v_pt=pad_v,
                for_verify=True,
            )
            if needed > cap * 1.02:
                preview = tf.text[:40].replace("\n", " ")
                issues.append(
                    f"slide {slide_idx}: text overflow "
                    f"(need ~{needed:.2f}\", cap ~{cap:.2f}\") [{preview}...]"
                )
    return issues


class SlideBuilder:
    def __init__(self, prs: Presentation) -> None:
        self.prs = prs
        self._slide_count = 0
        self._anim_seq = 0

    def _reset_anim(self) -> None:
        self._anim_seq = 0

    def _tag_anim(self, shape) -> None:
        self._anim_seq += 1
        shape.name = f"anim_{self._anim_seq:03d}"

    def _blank_slide(self):
        self._reset_anim()
        self._slide_count += 1
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        bg = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = WHITE
        bg.line.fill.background()
        return slide

    def _title_bar(self, slide, title: str) -> None:
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, HEADER_H
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY
        bar.line.fill.background()
        box = slide.shapes.add_textbox(MARGIN_L, Inches(0.18), CONTENT_W, Inches(0.55))
        tf = box.text_frame
        _configure_text_frame(tf, anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = FONT_UI
        p.font.size = FONT_TITLE
        p.font.bold = True
        p.font.color.rgb = WHITE

    def _accent_panel(
        self, slide, top: float, height: float, *, fill: RGBColor = PANEL_SOFT
    ) -> None:
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            MARGIN_L,
            Inches(top),
            CONTENT_W,
            Inches(height),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = fill
        panel.line.color.rgb = PANEL_BORDER
        panel.line.width = Pt(0.75)
        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            MARGIN_L,
            Inches(top),
            Inches(0.08),
            Inches(height),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = PRIMARY
        accent.line.fill.background()

    def _badge_pill(self, slide, badge: str, top: float) -> float:
        fill = AMBER if "Stage 4" in badge else ACCENT
        w = min(Inches(7.0), Inches(0.19 * len(badge) + 1.6))
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            MARGIN_L,
            Inches(top),
            w,
            Inches(0.52),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
        tf = shape.text_frame
        _configure_text_frame(tf, anchor=MSO_ANCHOR.MIDDLE)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = badge
        p.font.name = FONT_UI
        p.font.size = FONT_BADGE
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        return top + 0.62

    def _header(self, slide, title: str, *, badge: str | None = None) -> float:
        self._title_bar(slide, title)
        top = 1.05
        if badge:
            top = self._badge_pill(slide, badge, top)
        return top

    def _add_textbox(
        self,
        slide,
        top: float,
        height: float,
        *,
        anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    ):
        box = slide.shapes.add_textbox(
            _textbox_left(),
            Inches(top),
            _textbox_width(),
            Inches(height),
        )
        _configure_text_frame(box.text_frame, anchor=anchor)
        return box

    def _para(
        self,
        para,
        text: str,
        *,
        font_name: str = FONT_UI,
        font_size=FONT_BODY,
        bold: bool = False,
        color: RGBColor = BODY,
        space_after: int = 4,
        line_spacing: float = 1.12,
    ) -> None:
        para.text = text
        para.font.name = font_name
        para.font.size = font_size
        para.font.bold = bold
        para.font.color.rgb = color
        para.space_after = Pt(space_after)
        para.line_spacing = line_spacing

    def _bullets(
        self,
        slide,
        bullets: list[str],
        top: float,
        *,
        panel: bool = False,
        warn_panel: bool = False,
    ) -> None:
        content_top = top + 0.12
        height = BOTTOM_Y - content_top
        if panel:
            fill = WARNING_PANEL if warn_panel else PANEL_SOFT
            self._accent_panel(slide, content_top - 0.06, height + 0.06, fill=fill)

        box = self._add_textbox(slide, content_top, height)
        tf = box.text_frame
        section_keys = ("★", "活动", "练习", "观点表达", "论据支撑", "逻辑衔接", "动笔前")
        warn_keys = ("❌", "易错", "理由空泛", "语言平实", "逻辑生硬", "语气不当", "元素模糊")

        for i, bullet in enumerate(bullets):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            is_section = any(bullet.startswith(k) for k in section_keys)
            is_warn = any(k in bullet for k in warn_keys)
            text = bullet
            if not is_section and not text.startswith(("•", "→")):
                text = f"• {bullet}"
            self._para(
                para,
                text,
                font_size=FONT_SECTION if is_section else FONT_BODY,
                bold=is_section,
                color=WARNING if is_warn else BODY,
                space_after=5,
            )

    def title_slide(self, title: str, subtitle: str, body_lines: list[str]) -> None:
        slide = self._blank_slide()
        cover_h = Inches(2.35)
        bar = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, cover_h
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = PRIMARY
        bar.line.fill.background()
        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), cover_h, SLIDE_W, Inches(0.05)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = ACCENT
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(MARGIN_L, Inches(0.55), CONTENT_W, Inches(0.7))
        _configure_text_frame(title_box.text_frame, anchor=MSO_ANCHOR.MIDDLE)
        self._para(
            title_box.text_frame.paragraphs[0],
            title,
            font_size=FONT_TITLE,
            bold=True,
            color=WHITE,
        )
        subtitle_box = slide.shapes.add_textbox(MARGIN_L, Inches(1.35), CONTENT_W, Inches(0.55))
        _configure_text_frame(subtitle_box.text_frame, anchor=MSO_ANCHOR.TOP)
        self._para(
            subtitle_box.text_frame.paragraphs[0],
            subtitle,
            font_size=FONT_SECTION,
            color=RGBColor(0xE0, 0xE7, 0xFF),
        )
        panel_top = 2.28
        panel_h = BOTTOM_Y - panel_top - 0.06
        self._accent_panel(slide, panel_top, panel_h, fill=PANEL_SOFT)
        inner_top = panel_top + COVER_PANEL_PAD
        inner_h = panel_h - COVER_PANEL_PAD * 2
        body = self._add_textbox(slide, inner_top, inner_h)
        _configure_text_frame(body.text_frame, anchor=MSO_ANCHOR.TOP, pad_h=14, pad_v=10)
        for i, line in enumerate(body_lines):
            if not line.strip():
                continue
            self._para(
                body.text_frame.paragraphs[0] if i == 0 else body.text_frame.add_paragraph(),
                line,
                space_after=10,
                line_spacing=1.18,
            )

    def content_slide(
        self,
        title: str,
        bullets: list[str],
        *,
        badge: str | None = None,
        panel: bool = False,
        warn_panel: bool = False,
    ) -> None:
        slide = self._blank_slide()
        top = self._header(slide, title, badge=badge)
        self._bullets(slide, bullets, top, panel=panel or warn_panel, warn_panel=warn_panel)

    def peel_slide(self, title: str, points: list[dict]) -> None:
        slide = self._blank_slide()
        top = self._header(slide, title)
        card_top = top + 0.1
        card_h = BOTTOM_Y - card_top - 0.12
        card_w = 5.85
        card_colors = (PRIMARY, ACCENT)
        card_fills = (PANEL_BG, PANEL_ALT)
        labels = ("① 选择", "② 理由")

        for idx, point in enumerate(points[:2]):
            x = MARGIN_L.inches + idx * (card_w + 0.25)
            color = card_colors[idx]
            card = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(x),
                Inches(card_top),
                Inches(card_w),
                Inches(card_h),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = card_fills[idx]
            card.line.color.rgb = color
            card.line.width = Pt(2)

            hdr = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(card_top + 0.12), Inches(card_w - 0.3), Inches(0.48)
            )
            _configure_text_frame(hdr.text_frame, anchor=MSO_ANCHOR.MIDDLE)
            self._para(
                hdr.text_frame.paragraphs[0],
                labels[idx],
                font_size=FONT_SECTION,
                bold=True,
                color=color,
                space_after=0,
            )

            star = slide.shapes.add_textbox(
                Inches(x + 0.15), Inches(card_top + 0.58), Inches(card_w - 0.3), Inches(0.38)
            )
            _configure_text_frame(star.text_frame)
            self._para(
                star.text_frame.paragraphs[0],
                f"★ {point.get('label', '')}",
                font_size=FONT_BODY,
                bold=True,
                color=BODY,
                space_after=2,
            )

            body = slide.shapes.add_textbox(
                Inches(x + 0.15),
                Inches(card_top + 0.96),
                Inches(card_w - 0.3),
                Inches(card_h - 1.05),
            )
            _configure_text_frame(body.text_frame, pad_h=4, pad_v=4)
            tf = body.text_frame
            first = True

            def _add_line(text: str, *, english: bool = False, prefix: str = "") -> None:
                nonlocal first
                if not text:
                    return
                para = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                self._para(
                    para,
                    f"{prefix}{text}",
                    font_name=FONT_ESSAY if english else FONT_UI,
                    font_size=FONT_TABLE,
                    color=BODY,
                    space_after=6,
                    line_spacing=1.12,
                )

            _add_line(point.get("p", ""), english=True, prefix="P · ")
            e_items = point.get("e_items") or []
            if e_items:
                for item in e_items:
                    _add_line(item, english=("：" not in item[:8]))
            elif point.get("e"):
                _add_line(point["e"], english=False, prefix="E · ")
            _add_line(point.get("l", ""), english=True, prefix="L · ")

    def essay_slide(
        self,
        title: str,
        essay_text: str,
        annotation: str,
        *,
        badge: str | None = None,
    ) -> None:
        slide = self._blank_slide()
        top = self._header(slide, title, badge=badge)
        paragraphs = classroom_body_paragraphs(essay_text)
        line_spacing, para_space_pt, indent_spaces = essay_layout_for_length(paragraphs)

        content_top = top + 0.06
        annotation_h = 0.34 if annotation else 0.0
        content_h = BOTTOM_Y - content_top - annotation_h - 0.04

        box = slide.shapes.add_textbox(
            MARGIN_L,
            Inches(content_top),
            CONTENT_W,
            Inches(content_h),
        )
        tf = box.text_frame
        _configure_text_frame(tf, anchor=MSO_ANCHOR.TOP, pad_h=10, pad_v=6)

        first = True
        for pi, block in enumerate(paragraphs):
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            para.text = (" " * indent_spaces) + block.strip()
            para.font.name = FONT_ESSAY
            para.font.size = FONT_ESSAY_SIZE
            para.line_spacing = line_spacing
            para.space_after = Pt(2)
            if pi > 0:
                para.space_before = Pt(para_space_pt)

        if annotation:
            note = slide.shapes.add_textbox(
                MARGIN_L,
                Inches(content_top + content_h + 0.02),
                CONTENT_W,
                Inches(annotation_h - 0.02),
            )
            _configure_text_frame(note.text_frame, anchor=MSO_ANCHOR.TOP)
            self._para(
                note.text_frame.paragraphs[0],
                annotation,
                font_size=FONT_ANNOTATION,
                color=SECONDARY,
                space_after=0,
                line_spacing=1.0,
            )

    def table_slide(self, title: str, headers: list[str], rows: list[list[str]]) -> None:
        slide = self._blank_slide()
        top = self._header(slide, title)
        n_rows = len(rows) + 1
        max_table_h = BOTTOM_Y - top - 0.12
        table_h = min(max_table_h, 0.55 + 0.62 * n_rows)
        table = slide.shapes.add_table(
            n_rows, len(headers), MARGIN_L, Inches(top + 0.15), CONTENT_W, Inches(table_h)
        ).table
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY
            _configure_text_frame(cell.text_frame, anchor=MSO_ANCHOR.MIDDLE)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_UI
                p.font.size = FONT_TABLE
                p.font.bold = True
                p.font.color.rgb = WHITE
        for r, row in enumerate(rows, start=1):
            for c, val in enumerate(row):
                cell = table.cell(r, c)
                cell.text = val
                _configure_text_frame(cell.text_frame, anchor=MSO_ANCHOR.MIDDLE)
                for p in cell.text_frame.paragraphs:
                    p.font.name = FONT_UI
                    p.font.size = FONT_TABLE
                    p.font.color.rgb = BODY
                if r % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = PANEL_ALT

    def _style_table_header(self, table, headers: list[str]) -> None:
        for c, h in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = PRIMARY
            _configure_text_frame(cell.text_frame, anchor=MSO_ANCHOR.MIDDLE)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_UI
                p.font.size = FONT_TABLE
                p.font.bold = True
                p.font.color.rgb = WHITE

    def _style_table_data_row(
        self, table, row_idx: int, values: list[str], *, alt: bool = False
    ) -> None:
        for c, val in enumerate(values):
            cell = table.cell(row_idx, c)
            cell.text = val
            _configure_text_frame(cell.text_frame, anchor=MSO_ANCHOR.TOP)
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_ESSAY if c == 1 and len(values) == 3 else FONT_UI
                p.font.size = FONT_TABLE
                p.font.color.rgb = BODY
            if alt:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PANEL_ALT

    def _phrase_tier_note(self, tier: dict) -> str:
        level = tier.get("level", "")
        if level == "基础句":
            return ""
        note = tier.get("chinese", "") or ""
        hs = tier.get("high_score")
        if hs:
            note = f"{note}\n{hs}" if note else hs
        return note

    def phrase_table_slide(
        self, title: str, table: dict, *, badge: str | None = None
    ) -> list[int]:
        slide = self._blank_slide()
        top = self._header(slide, title, badge=badge)
        headers = ["层级", "英文句型", "说明"]
        col_fracs = [0.15, 0.55, 0.30]
        tiers = table.get("tiers", [])

        footer_h = 0.0
        if table.get("topic_note"):
            footer_h += _table_row_height([table["topic_note"]], [1.0]) + 0.12
        if table.get("fix_bad") or table.get("fix_good"):
            fix_text = " ".join(
                x for x in (table.get("fix_bad", ""), table.get("fix_good", "")) if x
            )
            footer_h += _table_row_height([fix_text], [1.0]) + 0.22

        content_top = top + 0.08
        row_heights = [0.52]
        row_values: list[list[str]] = []
        for tier in tiers:
            vals = [
                tier.get("level", ""),
                tier.get("english", ""),
                self._phrase_tier_note(tier),
            ]
            row_values.append(vals)
            row_heights.append(_table_row_height(vals, col_fracs))

        max_table_h = BOTTOM_Y - content_top - footer_h - 0.12
        total_h = sum(row_heights)
        if total_h > max_table_h and total_h > 0:
            scale = max_table_h / total_h
            row_heights = [max(0.40, h * scale) for h in row_heights]

        tbl_shape = slide.shapes.add_table(
            len(row_heights),
            3,
            MARGIN_L,
            Inches(content_top),
            CONTENT_W,
            Inches(sum(row_heights)),
        )
        tbl = tbl_shape.table
        col_w = CONTENT_W.inches
        for i, frac in enumerate(col_fracs):
            tbl.columns[i].width = Inches(col_w * frac)
        for ri, h in enumerate(row_heights):
            tbl.rows[ri].height = Inches(h)

        self._style_table_header(tbl, headers)
        anim_ids: list[int] = []
        for ri, vals in enumerate(row_values, start=1):
            self._style_table_data_row(tbl, ri, vals, alt=ri % 2 == 0)
        self._tag_anim(tbl_shape)
        anim_ids.append(self._anim_seq)

        cursor = content_top + sum(row_heights) + 0.10
        if table.get("topic_note"):
            note_h = _table_row_height([table["topic_note"]], [1.0])
            box = self._add_textbox(slide, cursor, note_h)
            self._para(
                box.text_frame.paragraphs[0],
                f"本题：{table['topic_note']}",
                font_size=FONT_BODY,
                color=BODY,
                line_spacing=1.12,
            )
            self._tag_anim(box)
            anim_ids.append(self._anim_seq)
            cursor += note_h + 0.10

        if table.get("fix_bad") or table.get("fix_good"):
            fix_lines = [x for x in (table.get("fix_bad"), table.get("fix_good")) if x]
            fix_h = _table_row_height(["\n".join(fix_lines)], [1.0]) + 0.16
            self._accent_panel(slide, cursor, fix_h, fill=PANEL_SOFT)
            inner = cursor + 0.08
            if table.get("fix_bad"):
                bad_h = _table_row_height([table["fix_bad"]], [1.0])
                bad_box = self._add_textbox(slide, inner, bad_h)
                self._para(
                    bad_box.text_frame.paragraphs[0],
                    table["fix_bad"],
                    color=WARNING,
                    line_spacing=1.10,
                )
                self._tag_anim(bad_box)
                anim_ids.append(self._anim_seq)
                inner += bad_h + 0.06
            if table.get("fix_good"):
                good_h = _table_row_height([table["fix_good"]], [1.0])
                good_box = self._add_textbox(slide, inner, good_h)
                self._para(
                    good_box.text_frame.paragraphs[0],
                    table["fix_good"],
                    color=ACCENT,
                    bold=True,
                    line_spacing=1.10,
                )
                self._tag_anim(good_box)
                anim_ids.append(self._anim_seq)
        return anim_ids

    def vocab_table_slide(
        self,
        title: str,
        tier: str,
        rows: list[dict],
        columns: list[str],
        *,
        badge: str | None = None,
    ) -> list[int]:
        slide = self._blank_slide()
        top = self._header(slide, title, badge=badge)
        header_map = {
            "english": "英文词块",
            "chinese": "中文释义",
            "example": "例句",
        }
        headers = [header_map[c] for c in columns]
        n_cols = len(columns)
        col_fracs = [1.0 / n_cols] * n_cols
        if n_cols == 3:
            col_fracs = [0.28, 0.22, 0.50]
        elif n_cols == 2:
            col_fracs = [0.38, 0.62]

        content_top = top + 0.08
        row_heights = [0.52]
        row_values: list[list[str]] = []
        for row in rows:
            vals = [row.get(c, "") for c in columns]
            row_values.append(vals)
            row_heights.append(_table_row_height(vals, col_fracs))

        max_table_h = BOTTOM_Y - content_top - 0.12
        total_h = sum(row_heights)
        if total_h > max_table_h and total_h > 0:
            scale = max_table_h / total_h
            row_heights = [max(0.40, h * scale) for h in row_heights]

        tbl_shape = slide.shapes.add_table(
            len(row_heights),
            n_cols,
            MARGIN_L,
            Inches(content_top),
            CONTENT_W,
            Inches(sum(row_heights)),
        )
        tbl = tbl_shape.table
        col_w = CONTENT_W.inches
        for i, frac in enumerate(col_fracs):
            tbl.columns[i].width = Inches(col_w * frac)
        for ri, h in enumerate(row_heights):
            tbl.rows[ri].height = Inches(h)

        self._style_table_header(tbl, headers)
        for ri, vals in enumerate(row_values, start=1):
            self._style_table_data_row(tbl, ri, vals, alt=ri % 2 == 0)
        self._tag_anim(tbl_shape)
        return [self._anim_seq]


def _bullet_weight(text: str) -> float:
    stripped = text.strip()
    if not stripped:
        return 0.35
    return max(1, len(stripped) // _CHARS_PER_LINE_BODY) * 1.0


def _split_bullets(bullets: list[str], max_weight: float) -> list[list[str]]:
    chunks: list[list[str]] = []
    current: list[str] = []
    weight = 0.0
    for bullet in bullets:
        w = _bullet_weight(bullet)
        if current and (weight + w > max_weight or len(current) >= MAX_BULLETS_PER_SLIDE):
            chunks.append(current)
            current = [bullet]
            weight = w
        else:
            current.append(bullet)
            weight += w
    if current:
        chunks.append(current)
    return chunks or [[]]


def expand_slide_specs(slides: list[dict]) -> list[dict]:
    expanded: list[dict] = []
    for spec in slides:
        if spec["type"] not in ("content",):
            expanded.append(spec)
            continue
        chunks = _split_bullets(spec["bullets"], MAX_BODY_WEIGHT)
        for idx, chunk in enumerate(chunks):
            new_spec = dict(spec)
            new_spec["bullets"] = chunk
            if len(chunks) > 1:
                new_spec["title"] = f"{spec['title']}（{idx + 1}/{len(chunks)}）"
                if idx > 0:
                    new_spec["badge"] = None
            expanded.append(new_spec)
    return expanded


def build_mental_health_deck() -> list[dict]:
    basic = (
        "Dear James,\n\n"
        "Glad to hear about your poster designs for the mental health week! "
        "I've looked at both, and I want to share my thoughts.\n\n"
        "I'd go with Poster 1. The cracked heart with a smile and the text "
        "\"It's okay not to be okay\" really stand out.\n\n"
        "The crack shows that everyone has tough times, and the smile makes it feel hopeful.\n\n"
        "Hope this helps! Good luck with the contest!"
    )
    high_a = (
        "Dear James,\n\n"
        "I was so excited to see your poster designs for the mental health week—they're both creative! "
        "I've been thinking about which one to pick, and I want to share my take.\n\n"
        "Personally, I'd choose Poster 1. The cracked heart with a smile instantly resonated with me. "
        "When I saw the crack, I thought of my own moments of vulnerability, like when I felt overwhelmed "
        "last semester.\n\n"
        "But the smile? It made me feel seen—like it's normal to have ups and downs, and that's okay. "
        "The text \"It's okay not to be okay\" wraps it all up perfectly, making the message "
        "feel personal and comforting.\n\n"
        "Anyway, I hope this helps! You've got great taste—good luck with the contest!"
    )
    high_b = (
        "Dear James,\n\n"
        "Thanks for sharing your poster designs for the mental health week! "
        "I've analyzed both, and I'd like to give my opinion.\n\n"
        "I recommend Poster 1. Firstly, the cracked heart symbolizes emotional vulnerability, which is a "
        "universal experience—everyone faces struggles, and acknowledging this is the first step to healing. "
        "Secondly, the smile conveys acceptance, a key part of mental well-being, and the text "
        "\"It's okay not to be okay\" reinforces this message, creating a cohesive theme. Unlike Poster 2, "
        "which focuses on growth (watering the heart), Poster 1 directly addresses the reality of emotional "
        "pain, making its message more immediate and impactful.\n\n"
        "Overall, Poster 1's design is more aligned with the core theme of mental health.\n\n"
        "Hope this helps! Good luck!"
    )
    return [
        {
            "type": "title",
            "title": "高考英语应用文 · 观点理由类",
            "subtitle": "心理健康周海报选题",
            "body": [
                "【题目】假如你是李华。交换生朋友 James 参加你校心理健康活动周"
                "海报设计大赛，创作了两个版本，向你征求意见。",
                "请你给他回复邮件：（1）你的选择；（2）说明理由。",
                "Poster 1：双手托着带裂痕但微笑的心 — \"It's okay not to be okay.\"",
                "Poster 2：浇水壶浇灌带心形叶子的植物 — \"Water your heart\"",
            ],
        },
        {
            "type": "content",
            "title": "审题 · 题型与动笔自检",
            "bullets": [
                "题目类型：观点理由类（题干要求「你的选择」+「说明理由」）",
                "动笔前自检五问：",
                "语气 — 朋友间建议 vs 正式报告？",
                "结构 — 主体段有具体理由还是只有「我觉得好」？",
                "逻辑 — 是否既说明「为什么选这个」，也解释「为什么不选另一个」？",
            ],
        },
        {
            "type": "content",
            "title": "审题 · 题型与动笔自检（续）",
            "bullets": [
                "立意 — 是否写出心理健康深层意义（接纳情绪、自我关怀）？",
                "语言 — 是否用「裂痕」「浇水」「自我关怀」等关键表达，避开 very good / nice？",
            ],
        },
        {
            "type": "content",
            "title": "审题 · 三元审题与交际目的",
            "bullets": [
                "我是谁：交换生朋友李华",
                "写给谁：交换生朋友 James",
                "为了什么：在心理健康周海报设计大赛中选择一个版本并说明理由",
                "核心交际目的：帮助 James 选择更合适的版本并理解选择理由",
                "达成标准：James 能清楚知道选择，觉得理由合理，愿意考虑",
            ],
        },
        {
            "type": "content",
            "title": "审题 · 体裁规范与能力维度",
            "bullets": [
                "体裁：邮件 | 时态：一般现在时 | 人称：第一人称 | 语气：友好、建议",
                "能力维度：观点表达与说理论证（表明立场→提供理由→辩证分析）",
                "底层思维路径：先确定选择 → 分析海报设计元素与心理健康主题关联 → 用具体细节支撑观点",
                "高分要点·观点具体性：说明哪个海报更符合心理健康周主题",
                "高分要点·论据支撑：裂痕象征真实情绪、双手表支持；浇水壶/植物象征主动照顾内心",
            ],
        },
        {
            "type": "content",
            "title": "审题 · 构思维度与段落规划",
            "bullets": [
                "维度1·设计元素分析：裂痕→真实情绪；双手→支持；浇水壶→主动关怀；植物→成长",
                "维度2·主题契合度：Poster 1「It's okay not to be okay.」直接回应情绪困扰",
                "Poster 2「Water your heart」强调主动行动",
                "开头段：问候 + 表明写作目的",
                "主体段：要点[1]选择（~30词）+ 要点[2]理由（~40词）",
            ],
        },
        {
            "type": "content",
            "title": "审题 · 一句大实话",
            "bullets": [
                "💡 最危险陷阱：只说「我觉得哪个好」而不具体分析设计元素与心理健康主题关联",
                "高分关键：用具体细节支撑选择",
                "Poster 1 裂痕象征真实情绪，微笑表明不完美也没关系",
                "Poster 2 浇水象征自我关怀，植物象征成长",
                "应当写成：朋友间建议；理由结合海报设计元素；结尾鼓励祝福",
            ],
        },
        {
            "type": "content",
            "title": "动笔易错",
            "badge": "Stage 4 · 审题后",
            "warn_panel": True,
            "panel": True,
            "bullets": [
                "易错1·理由空泛：The poster is good because it's about mental health.",
                "→ 用 cracked heart / emotional vulnerability 等词块绑定主题",
                "易错2·语气不当：I suggest that you choose Poster 1.（过于正式）",
                "→ 改用 I'd go with Poster 1. / Personally, I'm leaning towards Poster 1.",
                "易错3·元素模糊：The poster has a heart. → 改用 the cracked heart with a smile / the watering can and heart-shaped leaves",
            ],
        },
        {
            "type": "peel",
            "title": "PEEL 写作骨架",
            "points": [
                {
                    "label": "Point 1 选择",
                    "p": "I'd go with Poster 1.",
                    "e_items": [
                        "具体化：the cracked heart with a smile and the text 'It's okay not to be okay'",
                        "感受：it immediately caught my eye",
                    ],
                    "l": "Here's why I think so.",
                },
                {
                    "label": "Point 2 说明理由",
                    "p": (
                        "The cracked heart with a smile and the text perfectly capture "
                        "the essence of mental health—acknowledging vulnerability while "
                        "embracing acceptance."
                    ),
                    "e_items": [
                        "具体化：裂痕代表挣扎，微笑表明不完美也没关系",
                        "影响：this combination makes the message relatable and comforting",
                    ],
                    "l": (
                        "Overall, Poster 1's design feels more authentic and "
                        "emotionally resonant."
                    ),
                },
            ],
        },
        {
            "type": "essay",
            "title": "基础版范文 · 9分档",
            "badge": "基础版 · 9分档",
            "essay_text": basic,
            "annotation": "（110 words · 内容齐全、语言平实 · 中间段约60词）",
        },
        {
            "type": "essay",
            "title": "高分版 A · 情感共鸣型",
            "badge": "高分版 A · 12–14分档",
            "essay_text": high_a,
            "annotation": "（118 words · 个人经历增强温度 · 定语从句/宾语从句）",
        },
        {
            "type": "essay",
            "title": "高分版 B · 逻辑思辨型",
            "badge": "高分版 B · 12–14分档",
            "essay_text": high_b,
            "annotation": "（122 words · Firstly/Secondly 逻辑链 · Unlike 比较结构）",
        },
        {
            "type": "table",
            "title": "三版对比要点",
            "headers": ["维度", "基础版", "高分版 A", "高分版 B"],
            "rows": [
                ["句式", "简单句为主", "定语/状语/宾语从句", "因果链、Unlike 比较"],
                ["词汇", "stand out, hopeful", "resonated, vulnerability", "symbolizes, cohesive"],
                ["衔接", "and, but", "情感过渡 when I saw...", "Firstly, Secondly, Unlike"],
                ["风格", "平实建议", "情感共鸣，人际温度", "逻辑思辨，严谨分析"],
                ["适合", "基础薄弱", "擅长情感表达", "擅长逻辑分析"],
            ],
        },
        {
            "type": "content",
            "title": "讲评活动 · 元素与逻辑",
            "badge": "Stage 4 · 范文后",
            "bullets": [
                "活动1·元素-主题关联拆解（10min）",
                "展示两张海报 → 圈关键元素 → 追问「与心理健康有何关联？」",
                "学生用词块造句：The cracked heart symbolizes emotional vulnerability.",
                "活动2·逻辑链模仿（15min）",
                "拆解高分版 B：Firstly 元素象征 → Secondly 主题强化",
            ],
        },
        {
            "type": "content",
            "title": "讲评活动 · 风格升级",
            "bullets": [
                "学生选海报，用 Firstly/Secondly 写两个理由",
                "活动3·风格选择与片段升级（10min）",
                "对比情感共鸣型 vs 逻辑思辨型，同伴互评语言层次",
            ],
        },
        {
            "type": "content",
            "title": "功能句型 · 观点表达",
            "bullets": [
                "基础：I prefer Poster 1.",
                "进阶：I'm leaning towards Poster 1, as its design stands out.",
                "高级：Among the two, Poster 1 strikes me as more compelling due to its emotional resonance.",
                "改一句：I think Poster 1 is good. ❌",
                "→ Personally, I'm leaning towards Poster 1, as its design really catches my eye. ✅",
            ],
        },
        {
            "type": "content",
            "title": "功能句型 · 论据支撑",
            "bullets": [
                "基础：Poster 1 is good because it has a cracked heart.",
                "进阶：Poster 1 resonates with me because the cracked yet smiling heart perfectly conveys the message of acceptance.",
                "高级：The visual metaphor of a cracked heart being held gently speaks volumes about the importance of support during tough times.",
                "改一句：The poster is good because it has a heart. ❌",
                "→ Poster 1 resonates with me because the cracked yet smiling heart perfectly conveys the message that it's normal to struggle but still find hope. ✅",
            ],
        },
        {
            "type": "content",
            "title": "功能句型 · 逻辑衔接",
            "bullets": [
                "基础：First, Poster 1 has a heart. Second, the text is good.",
                "进阶：Firstly, the visual of hands holding a cracked heart is powerful; secondly, the slogan directly addresses the core theme.",
                "高级：On one hand, the imagery of a cracked heart being supported highlights emotional vulnerability; on the other hand, the slogan normalizes struggles, making the message more relatable.",
            ],
        },
        {
            "type": "content",
            "title": "功能句型 · 逻辑衔接（续）",
            "bullets": [
                "改一句：I like Poster 1. It has a heart. The text is good. ❌",
                "→ Firstly, the visual of hands holding a cracked heart is powerful; secondly, the slogan \"It's okay not to be okay\" directly addresses the core theme of mental health acceptance. ✅",
            ],
        },
        {
            "type": "content",
            "title": "话题词块 · 观点表达",
            "bullets": [
                "必备：prefer / think... is better / have a preference for",
                "进阶：lean towards / find... more compelling / be inclined to choose",
                "亮点：strike me as more impactful / be drawn to / favor... over...",
                "例句：I'm leaning towards Poster 1, as its message is more direct.",
            ],
        },
        {
            "type": "content",
            "title": "话题词块 · 设计元素",
            "bullets": [
                "必备：cracked heart / watering can / heart-shaped leaves",
                "进阶：visual of hands holding / smiling heart / metaphorical imagery",
                "亮点：symbolic representation / visual metaphor / raw honesty",
                "例句：The visual of hands holding a cracked heart conveys support.",
            ],
        },
        {
            "type": "content",
            "title": "话题词块 · 心理健康主题",
            "bullets": [
                "必备：mental health / acceptance / self-care",
                "进阶：emotional vulnerability / normalizing struggles / fostering resilience",
                "亮点：emotional resonance / relatable message / nurturing mental well-being",
                "例句：The cracked heart reflects emotional vulnerability, which is a key part of mental health.",
            ],
        },
        {
            "type": "content",
            "title": "当堂迁移 · 片段升级",
            "badge": "Stage 4 · 句型后",
            "bullets": [
                "练习1·片段升级（选 Poster 2）：",
                "基础：I choose Poster 2. It has a watering can and a plant.",
                "升级：I'd go with Poster 2. The watering can and the plant with heart-shaped leaves really stand out. The watering can represents self-care, which is key to nurturing mental well-being, and the plant symbolizes growth.",
            ],
        },
        {
            "type": "content",
            "title": "当堂迁移 · 完整写作",
            "bullets": [
                "练习2·完整写作迁移：",
                "新题：Lucy 校园环保海报（版本A地球被塑料袋包围 vs 版本B双手种树）",
                "考查：论据支撑（结合设计元素与环保主题）",
                "提示：设计元素词块 + Firstly/Secondly + 选情感共鸣或逻辑思辨风格",
            ],
        },
        {
            "type": "content",
            "title": "易错深化 · 动笔改错",
            "badge": "Stage 4 · 融入",
            "warn_panel": True,
            "panel": True,
            "bullets": [
                "1. 理由空泛 → 用设计元素词块绑定心理健康主题",
                "2. 语言平实 → 用进阶句型拓展元素意义（resonates / symbolizes）",
                "3. 逻辑生硬 → 用 Firstly/Secondly 层次推进，避免流水账",
                "4. 语气不当 → 朋友间用 I'd go with / Personally, I'm leaning towards",
                "5. 元素模糊 → 具体描述 the cracked heart with a smile",
            ],
        },
        {
            "type": "content",
            "title": "课堂小结",
            "bullets": [
                "审题：观点理由类 = 选择 + 理由，结合设计元素与主题",
                "范文：PEEL 骨架 → 三版对比 → 选情感共鸣或逻辑思辨风格",
                "句型：观点表达 + 论据支撑 + 逻辑衔接 + 话题词块",
                "迁移：用同样方法分析新场景（环保海报选题）",
                "下节课：同伴互评 + 完整邮件写作",
            ],
        },
    ]


def _clean_pptx_dir(output: Path) -> None:
    """Remove Office lock files only — keep V1/V2 and other lesson decks side by side."""
    out_dir = output.parent
    if not out_dir.is_dir():
        return
    for lock in out_dir.glob("~$*.pptx"):
        try:
            lock.unlink()
        except OSError:
            pass


def _enforce_word_wrap(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell.text_frame.word_wrap = True
            if shape.has_text_frame:
                shape.text_frame.word_wrap = True


def render_deck(slides: list[dict], output: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    builder = SlideBuilder(prs)
    for spec in slides:
        kind = spec["type"]
        if kind == "title":
            builder.title_slide(spec["title"], spec["subtitle"], spec["body"])
        elif kind == "essay":
            builder.essay_slide(
                spec["title"],
                spec["essay_text"],
                spec.get("annotation", ""),
                badge=spec.get("badge"),
            )
        elif kind == "table":
            builder.table_slide(spec["title"], spec["headers"], spec["rows"])
        elif kind == "peel":
            builder.peel_slide(spec["title"], spec["points"])
        elif kind == "phrase_table":
            builder.phrase_table_slide(spec["title"], spec["table"], badge=spec.get("badge"))
        elif kind == "vocab_table":
            builder.vocab_table_slide(
                spec["title"],
                spec.get("tier", ""),
                spec["rows"],
                spec["columns"],
                badge=spec.get("badge"),
            )
        else:
            builder.content_slide(
                spec["title"],
                spec["bullets"],
                badge=spec.get("badge"),
                panel=spec.get("panel", False),
                warn_panel=spec.get("warn_panel", False) or "易错" in spec["title"],
            )
    output.parent.mkdir(parents=True, exist_ok=True)
    _enforce_word_wrap(prs)
    _clean_pptx_dir(output)
    prs.save(output)
    return output


def build_deck_with_stage3(
    stage3_path: Path,
    deck_plan_path: Path | None = None,
    *,
    vocab_max_rows: int = 6,
) -> list[dict]:
    import json

    stage3_data = json.loads(stage3_path.read_text(encoding="utf-8"))
    plan = load_deck_plan(deck_plan_path, stage3_data, vocab_max_rows=vocab_max_rows)
    stage3_specs = stage3_specs_from_plan(stage3_data, plan)
    deck = build_mental_health_deck()
    return expand_slide_specs(replace_stage3_in_deck(deck, stage3_specs))


def _collect_font_sizes(prs: Presentation) -> tuple[set[int], list[int]]:
    seen: set[int] = set()
    violations: list[int] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            if para.font.size:
                                pt = int(round(para.font.size.pt))
                                seen.add(pt)
                                if pt < MIN_FONT_PT:
                                    violations.append(pt)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.font.size:
                        pt = int(round(para.font.size.pt))
                        seen.add(pt)
                        if pt < MIN_FONT_PT:
                            violations.append(pt)
    return seen, violations


_PAGE_NUM_RE = re.compile(r"^\s*\d+\s*/\s*\d+\s*$")


def _find_page_numbers(prs: Presentation) -> list[str]:
    hits: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                if _PAGE_NUM_RE.match(para.text.strip()):
                    hits.append(f"slide {slide_idx}: {para.text.strip()}")
    return hits


_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _scan_pptx_xml_pagination(path: Path) -> list[str]:
    hits: list[str] = []
    with zipfile.ZipFile(path, "r") as zf:
        for name in sorted(zf.namelist()):
            if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(zf.read(name))
            for t_elem in root.iter(f"{_A_NS}t"):
                text = (t_elem.text or "").strip()
                if _PAGE_NUM_RE.match(text):
                    hits.append(f"{name}: {text}")
    return hits


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Generate 应用文 classroom PPTX")
    parser.add_argument(
        "-o",
        "--output",
        type=Path,
        default=Path(r"d:\Downloads\ppt-work\mental_health_classroom_V1.pptx"),
    )
    parser.add_argument(
        "--stage3",
        type=Path,
        default=None,
        help="stage3.json path (enables Stage3 table replacement)",
    )
    parser.add_argument(
        "--deck-plan",
        type=Path,
        default=None,
        help="classroom_deck.json slide plan (optional)",
    )
    parser.add_argument(
        "--no-anim",
        action="store_true",
        help="skip WPS on-click animation injection",
    )
    args = parser.parse_args(argv)
    if args.stage3:
        slides = build_deck_with_stage3(args.stage3.expanduser().resolve(), args.deck_plan)
    else:
        slides = expand_slide_specs(build_mental_health_deck())
    path = render_deck(slides, args.output)
    if not args.no_anim:
        apply_click_reveal(path)
    prs = Presentation(path)
    seen, violations = _collect_font_sizes(prs)
    overflow = verify_text_fit(prs)
    page_nums = _find_page_numbers(prs)
    xml_page_nums = _scan_pptx_xml_pagination(path)
    essays = [s for s in slides if s["type"] == "essay"]

    print(f"Saved: {path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Font sizes used: {sorted(seen)}")
    print(f"Min font: {min(seen) if seen else 'n/a'}")
    print(f"All >= {MIN_FONT_PT}pt: {'yes' if not violations else 'NO ' + str(sorted(set(violations)))}")
    print(f"Page numbers (shapes): {'none' if not page_nums else 'FOUND ' + str(page_nums)}")
    print(f"Page numbers (XML): {'none' if not xml_page_nums else 'FOUND ' + str(xml_page_nums)}")
    print(f"Text overflow check: {'pass' if not overflow else 'FAIL'}")
    for issue in overflow:
        print(f"  - {issue}")
    print(f"Essay slides: {len(essays)} (each full text, no Part 1/2)")
    for s in essays:
        has_dear = "Dear James" in s["essay_text"]
        has_end = "Good luck" in s["essay_text"] or "Hope this helps" in s["essay_text"]
        print(f"  {s['title']}: Dear James={has_dear}, sign-off={has_end}")
    ok = not violations and not overflow and not page_nums and not xml_page_nums
    return 0 if ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
