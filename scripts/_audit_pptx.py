#!/usr/bin/env python
"""One-off audit for mental_health_classroom.pptx — write report to file."""
from __future__ import annotations

import json
import sys
from pathlib import Path

_PROJECT = Path(__file__).resolve().parents[1]
if str(_PROJECT) not in sys.path:
    sys.path.insert(0, str(_PROJECT))

from pptx import Presentation

from scripts.classroom_script import compile_classroom_script, load_script_template
from scripts.deck_plan import deck_plan_from_stage3, refine_deck_plan, stage3_specs_from_plan
from scripts.generate_classroom_pptx import _collect_font_sizes
from scripts.wps_layout_verify import verify_deck_layout
from scripts.ppt_layout_fit import (
    expand_content_slides,
    expand_essay_slides,
    expand_peel_slides,
    expand_title_slides,
)

PPTX = Path(r"D:\Downloads\ppt-work\mental_health_classroom.pptx")
OUT = Path(r"D:\Downloads\ppt-work\audit-report.txt")
WORKDIR = Path(r"D:\Downloads\ppt-work")


def _slide_texts(slide) -> list[tuple[str, float, float, float, float, int | None]]:
    """Extract text shapes: (preview, top, left, w, h, font_pt_min)."""
    items: list[tuple[str, float, float, float, float, int | None]] = []
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        text = tf.text.strip()
        if not text:
            continue
        tops = []
        for p in tf.paragraphs:
            for r in p.runs:
                if r.font.size:
                    tops.append(int(r.font.size.pt))
        min_pt = min(tops) if tops else None
        preview = text.replace("\n", " ")[:80]
        items.append(
            (
                preview,
                sh.top.inches if sh.top else 0,
                sh.left.inches if sh.left else 0,
                sh.width.inches if sh.width else 0,
                sh.height.inches if sh.height else 0,
                min_pt,
            )
        )
    return items


def _overlap_pairs(items: list[tuple]) -> list[str]:
    pairs: list[str] = []
    for i, a in enumerate(items):
        _, at, al, aw, ah, _ = a
        for j, b in enumerate(items[i + 1 :], i + 1):
            _, bt, bl, bw, bh, _ = b
            if at + ah > bt + 0.02 and bt + bh > at + 0.02:
                if al + aw > bl + 0.02 and bl + bw > al + 0.02:
                    pairs.append(f"shape {i+1} vs {j+1}: y=[{at:.2f}-{at+ah:.2f}] vs [{bt:.2f}-{bt+bh:.2f}]")
    return pairs


def main() -> None:
    lines: list[str] = []
    prs = Presentation(str(PPTX))
    layout = verify_deck_layout(prs)
    overflow = layout["pass1_issues"]
    wps_issues = layout["pass2_issues"]
    seen, violations = _collect_font_sizes(prs)

    lines.append("=" * 60)
    lines.append("PPTX AUDIT: mental_health_classroom.pptx")
    lines.append("=" * 60)
    lines.append(f"Slide count: {len(prs.slides)}")
    lines.append(f"Font sizes used: {sorted(seen)}")
    lines.append(f"Below 26pt violations: {sorted(set(violations)) if violations else 'none'}")
    lines.append(f"Overflow issues (pass1): {len(overflow)}")
    wps_report = layout.get("wps_report")
    if wps_report:
        lines.append(f"WPS critical: {len(wps_report.critical_issues)}")
        lines.append(f"WPS warning: {len(wps_report.warning_issues)}")
        lines.append(f"WPS cosmetic ignored: {len(wps_report.cosmetic_issues)}")
        lines.append(f"teach_ready: {wps_report.is_teach_ready}")
        lines.append(f"risk_score: {wps_report.risk_score:.3f}")
    lines.append(f"WPS_RISK_OVERFLOW (critical only): {len(wps_issues)}")
    lines.append("")

    # Per-slide dump
    lines.append("--- PER-SLIDE INSPECTION ---")
    small_font_slides: list[int] = []
    overlap_slides: list[int] = []
    for idx, slide in enumerate(prs.slides, 1):
        items = _slide_texts(slide)
        header = items[0][0][:50] if items else "(empty)"
        lines.append(f"\n[Slide {idx}] shapes_with_text={len(items)} | first: {header}")
        for k, (prev, t, l, w, h, pt) in enumerate(items, 1):
            flag = ""
            if pt is not None and pt < 26:
                flag = " **FONT<26**"
            lines.append(f"  {k}. y={t:.2f} h={h:.2f} w={w:.2f} pt={pt}{flag} | {prev}")
        ov = _overlap_pairs(items)
        if ov:
            overlap_slides.append(idx)
            for o in ov:
                lines.append(f"  OVERLAP: {o}")

    lines.append("\n--- WPS SAFETY DETAILS (first 40) ---")
    for o in wps_issues[:40]:
        lines.append(o)
    if len(wps_issues) > 40:
        lines.append(f"... +{len(wps_issues) - 40} more")

    lines.append("\n--- OVERFLOW DETAILS (pass1, first 40) ---")
    for o in overflow[:40]:
        lines.append(o)
    if len(overflow) > 40:
        lines.append(f"... +{len(overflow) - 40} more")

    # Script plan vs actual
    lines.append("\n--- SCRIPT PLAN (expanded) ---")
    export = json.loads((WORKDIR / "export-data.json").read_text(encoding="utf-8"))
    stage3 = json.loads((WORKDIR / "stage3.json").read_text(encoding="utf-8"))
    specs = stage3_specs_from_plan(stage3, refine_deck_plan(stage3, deck_plan_from_stage3(stage3)))
    compiled = compile_classroom_script(
        load_script_template("dual_poster_opinion"), export, stage3, stage3_specs=specs
    )
    expanded = expand_content_slides(
        expand_peel_slides(expand_essay_slides(expand_title_slides(compiled)))
    )
    lines.append(f"Compiled specs: {len(compiled)}")
    lines.append(f"After expand_*: {len(expanded)}")
    lines.append(f"Actual pptx: {len(prs.slides)}")
    if len(expanded) != len(prs.slides):
        lines.append(f"** MISMATCH: expand plan {len(expanded)} != pptx {len(prs.slides)}")

    for i, s in enumerate(expanded, 1):
        lines.append(f"  {i:2d} {s.get('type','?'):14s} {(s.get('title') or '')[:45]}")

    # Structural issues checklist
    lines.append("\n--- STRUCTURAL ISSUES (automated) ---")
    issues: list[str] = []

    # Cover should have empty body in spec
    if compiled and compiled[0].get("type") == "title":
        if compiled[0].get("body"):
            issues.append("Cover spec still has body text (should be empty)")
        if compiled[0].get("poster_lines"):
            issues.append("Cover spec still has poster_lines")

    # Poster count
    poster_specs = [s for s in compiled if s.get("title") == "海报示意"]
    poster_lines = sum(len(s.get("bullets") or []) for s in poster_specs)
    q = export.get("question") or ""
    if "[图" in q and poster_lines == 0:
        issues.append("Question has [图] but no poster slides compiled")

    # thinking_path + thinking_formula duplicate source?
    seq = load_script_template("dual_poster_opinion")["sequence"]
    sources = [item.get("bind", {}).get("source") for item in seq if item.get("bind", {}).get("source")]
    if len(sources) != len(set(sources)):
        issues.append(f"Template has duplicate sources: {sources}")

    # Stage3 phrase table count
    phrase_count = sum(1 for s in expanded if s.get("type") == "phrase_table")
    if phrase_count > 6:
        issues.append(f"Too many phrase_table slides ({phrase_count}) — may feel repetitive")

    vocab_count = sum(1 for s in expanded if s.get("type") == "vocab_table")
    if vocab_count > 6:
        issues.append(f"Too many vocab_table slides ({vocab_count}) — 3 fields × 3 tiers = 9 pages")

    # Essay annotation missing
    essays = [s for s in compiled if s.get("type") == "essay"]
    for e in essays:
        if not (e.get("annotation") or "").strip():
            issues.append(f"Essay '{e.get('title')}' has empty annotation in spec")

    # Migration content
    mig = next((s for s in compiled if "迁移" in (s.get("title") or "")), None)
    if mig:
        bullets = " ".join(mig.get("bullets") or [])
        if "Lucy" in bullets or "环保" in bullets:
            issues.append("Migration slide contains wrong topic (Lucy/环保)")
        if "James" not in bullets and "心理健康" not in bullets:
            issues.append("Migration slide may not anchor to James/心理健康")

    # Duplicate titles in expanded
    titles = [s.get("title") for s in expanded]
    from collections import Counter

    dup = {t: c for t, c in Counter(titles).items() if c > 1 and t}
    if dup:
        issues.append(f"Duplicate slide titles in plan: {dup}")

    if violations:
        issues.append(f"Renderer used 24pt on slides (violations count={len(violations)})")
    if len(wps_issues) > 50:
        issues.append(f"Mass WPS risk ({len(wps_issues)} shapes) — WPS layout mismatch")
    if len(overflow) > 50:
        issues.append(f"Mass overflow ({len(overflow)} shapes) — layout engine vs WPS mismatch")

    if overlap_slides:
        issues.append(f"Potential text box overlap on slides: {overlap_slides[:15]}")

    # Poster lines content check from export
    from scripts.architecture_v1 import _poster_lines, _question_lines

    posters = _poster_lines(q)
    stems = _question_lines(q)
    lines.append(f"\nQuestion stem lines: {len(stems)}")
    for s in stems:
        lines.append(f"  stem: {s[:100]}")
    lines.append(f"Poster lines parsed: {len(posters)}")
    for p in posters:
        lines.append(f"  poster: {p[:100]}")

    for iss in issues:
        lines.append(f"ISSUE: {iss}")

    lines.append("\n--- END ---")
    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {OUT} ({len(lines)} lines)")


if __name__ == "__main__":
    main()
