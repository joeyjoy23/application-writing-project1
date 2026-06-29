"""Unified PPT card primitive tests."""

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from styles.design_tokens import Theme as G
from styles.ppt_card import add_unified_card


def test_add_unified_card_creates_card_and_accent_bar(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    card = add_unified_card(
        slide,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(1.2),
    )
    assert card.fill.fore_color.rgb == G.SURFACE
    assert card.line.color.rgb == G.BORDER
    assert len(slide.shapes) >= 2
    assert any(s.shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE for s in slide.shapes)
    out = tmp_path / "card.pptx"
    prs.save(out)
    assert out.stat().st_size > 1000
