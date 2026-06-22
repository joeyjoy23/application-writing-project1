"""Tests for V1 Stage3 table slides in generate_classroom_pptx."""

import json
import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation

from scripts.generate_classroom_pptx import (
    SlideBuilder,
    build_deck_with_stage3,
    render_deck,
)
from scripts.parse_stage3 import parse_stage3_file, write_stage3_json
from scripts.pptx_click_reveal import apply_click_reveal

FIXTURE = Path(__file__).parent / "fixtures" / "stage3_mental_health.md"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_ANIM_RE = re.compile(r"^anim_\d{3}$")


def _anim_names_from_pptx(path: Path) -> list[str]:
    names: list[str] = []
    with zipfile.ZipFile(path, "r") as zf:
        for name in sorted(zf.namelist()):
            if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(zf.read(name))
            for cnv in root.iter(f"{{{_P_NS}}}cNvPr"):
                n = cnv.get("name", "")
                if _ANIM_RE.match(n):
                    names.append(n)
    return names


def test_phrase_table_row_count(tmp_path):
    data = parse_stage3_file(FIXTURE)
    table = data["phrase_tables"][0]
    prs = Presentation()
    builder = SlideBuilder(prs)
    builder.phrase_table_slide("功能句型 · 观点表达", table)
    slide = prs.slides[0]
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 1
    assert len(tables[0].table.rows) == 4  # header + 3 tiers


def test_vocab_basic_two_columns(tmp_path):
    data = parse_stage3_file(FIXTURE)
    field = next(f for f in data["vocab_fields"] if "设计" in f["name"])
    basic = next(t for t in field["tiers"] if t["level"] == "必备级")
    prs = Presentation()
    builder = SlideBuilder(prs)
    builder.vocab_table_slide(
        "话题词块 · 设计 · 必备级",
        "必备级",
        basic["rows"][:3],
        ["english", "example"],
    )
    slide = prs.slides[0]
    hdr = next(s for s in slide.shapes if s.has_table)
    assert hdr.table.columns.__len__() == 2


def test_build_deck_with_stage3_replaces_bullets(tmp_path):
    stage3_json = tmp_path / "stage3.json"
    write_stage3_json(parse_stage3_file(FIXTURE), stage3_json)
    slides = build_deck_with_stage3(stage3_json)
    types = [s["type"] for s in slides]
    assert "phrase_table" in types
    assert "vocab_table" in types
    # old bullet stage3 slides replaced
    stage3_bullets = [
        s
        for s in slides
        if s.get("type") == "content" and s.get("title", "").startswith("功能句型")
    ]
    assert stage3_bullets == []


def test_render_deck_anim_shape_names(tmp_path):
    stage3_json = tmp_path / "stage3.json"
    out = tmp_path / "deck.pptx"
    write_stage3_json(parse_stage3_file(FIXTURE), stage3_json)
    slides = build_deck_with_stage3(stage3_json)
    render_deck(slides, out)
    apply_click_reveal(out)
    names = _anim_names_from_pptx(out)
    assert len(names) >= 10
    assert all(_ANIM_RE.match(n) for n in names)
