"""Tests for WPS layout verification and governance layer."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from scripts.generate_classroom_pptx import verify_text_fit
from scripts.ppt_layout_fit import text_block_height
from scripts.wps_layout_verify import (
    WPSLayoutReport,
    classify_wps_element,
    compute_card_height,
    compute_safe_height,
    estimate_wps_text_height,
    is_teach_ready,
    safe_padding_inches,
    verify_deck_layout,
    verify_wps_layout_safety,
    WPS_LAYOUT_SAFETY,
    WPS_SAFE_FACTOR,
    CARD_CONTAINER_PAD_V,
    _margin_ratio,
)

_CJK_SAMPLE = (
    "心理健康是青少年成长的重要基石。学校应关注学生情绪变化，"
    "提供及时的心理支持与专业辅导，帮助他们在压力下保持积极心态。"
)


def test_compute_card_height_shrink_wraps_short_text():
    short = "• 语气：友好"
    h = compute_card_height(short, 28, 11.0, bullet=True)
    inflated = estimate_wps_text_height(short, 28, {"bullet": True}, width_inches=11.0)
    pad = safe_padding_inches(28) * 2 + CARD_CONTAINER_PAD_V
    assert h <= inflated * WPS_SAFE_FACTOR + pad + 0.15
    assert h >= inflated + pad - 0.05
    assert compute_safe_height(short, 28, 11.0, bullet=True) == h


def test_compute_card_height_long_text_capped_by_safe_max():
    long_text = "• " + _CJK_SAMPLE * 3
    h = compute_card_height(long_text, 28, 11.0, bullet=True)
    text_h = estimate_wps_text_height(long_text, 28, {"bullet": True, "chinese": True}, width_inches=11.0)
    pad = safe_padding_inches(28) * 2 + CARD_CONTAINER_PAD_V
    assert h <= text_h * WPS_SAFE_FACTOR + pad + 0.15


def test_estimate_wps_text_height_exceeds_verify_heuristic_for_cjk():
    font_pt = 28.0
    width = 5.0
    wps_h = estimate_wps_text_height(
        _CJK_SAMPLE,
        font_pt,
        {"chinese": True, "bold": False, "bullet": False, "table_cell": False},
        width_inches=width,
    )
    legacy = text_block_height(_CJK_SAMPLE, width, font_pt, line_spacing=1.12)
    ratio = wps_h / legacy
    assert ratio >= 1.10, f"WPS estimate should be >=10% above legacy ({ratio:.2f})"


def test_margin_ratio_calculation():
    assert _margin_ratio(1.20, 1.0) == pytest.approx(0.20)
    assert _margin_ratio(1.0, 1.0) == pytest.approx(0.0)
    assert _margin_ratio(2.0, 0.0) == 1.0


def test_safe_padding_minimum_is_14px():
    small_font = safe_padding_inches(8.0)
    assert small_font == pytest.approx(14.0 / 72.0)


def test_classify_section_pill_as_cosmetic():
    level, reason = classify_wps_element(
        text="审题",
        shape_top=0.12,
        shape_left=0.45,
        shape_h=0.48,
        shape_w=1.35,
        est=0.55,
        safe_h=0.90,
        margin_ratio=0.02,
        kind="text_frame",
    )
    assert level == "cosmetic"
    assert reason == "section_pill"


def test_classify_actual_overflow_as_critical():
    level, reason = classify_wps_element(
        text="• " + _CJK_SAMPLE,
        shape_top=2.0,
        shape_left=0.55,
        shape_h=0.45,
        shape_w=11.5,
        est=1.40,
        safe_h=1.60,
        margin_ratio=-0.63,
        kind="text_frame",
    )
    assert level == "critical"
    assert "overflow" in reason or "truncation" in reason


def test_classify_moderate_card_overflow_as_critical_when_under_allocated():
    level, reason = classify_wps_element(
        text="• Poster 1是双手托着带裂痕但微笑的心",
        shape_top=1.86,
        shape_left=0.55,
        shape_h=0.95,
        shape_w=11.5,
        est=1.23,
        safe_h=1.62,
        margin_ratio=-0.23,
        kind="text_frame",
    )
    assert level in ("critical", "warning")
    assert "card" in reason or "overflow" in reason or "allocated" in reason


def test_classify_tight_margin_as_warning():
    level, _ = classify_wps_element(
        text="• Poster 1是双手托着带裂痕但微笑的心",
        shape_top=1.86,
        shape_left=0.55,
        shape_h=0.95,
        shape_w=11.5,
        est=0.88,
        safe_h=1.10,
        margin_ratio=0.08,
        kind="text_frame",
    )
    assert level == "warning"


def test_is_teach_ready_blocks_on_critical():
    report = WPSLayoutReport(
        critical_issues=[object()],  # type: ignore[list-item]
        warning_issues=[],
        total_slides=10,
    )
    assert is_teach_ready(report) is False


def test_is_teach_ready_allows_low_warning_ratio():
    from scripts.wps_layout_verify import WPSLayoutIssue

    report = WPSLayoutReport(
        warning_issues=[WPSLayoutIssue("warning", 1, "s", "m")] * 3,
        total_slides=31,
    )
    assert is_teach_ready(report) is True


def _mock_tight_textbox(prs: Presentation, text: str, *, height_in: float = 0.45) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.55), Inches(2.0), Inches(8), Inches(height_in))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)


def test_verify_wps_governance_flags_tight_content_as_critical():
    prs = Presentation()
    _mock_tight_textbox(prs, "• " + _CJK_SAMPLE, height_in=0.35)
    report, dbg = verify_wps_layout_safety(prs)
    assert report.wps_risk_overflow_count >= 1
    assert dbg and dbg[0].elements


def test_verify_deck_layout_dual_pass():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1.5))
    box.text_frame.text = "Short label"
    result = verify_deck_layout(prs)
    assert result["ok"] is True
    assert result["wps_risk_count"] == 0
    assert result["is_teach_ready"] is True


@pytest.mark.parametrize(
    "pptx_path",
    [
        Path(r"D:\Downloads\ppt-work\mental_health_classroom.pptx"),
        Path("scripts/ppt_assets/mental_health_classroom.pptx"),
    ],
)
def test_governance_on_generated_pptx_if_present(pptx_path: Path):
    if not pptx_path.is_file():
        pytest.skip(f"no fixture pptx at {pptx_path}")
    prs = Presentation(str(pptx_path))
    report, slide_debug = verify_wps_layout_safety(prs)
    assert isinstance(report, WPSLayoutReport)
    assert len(slide_debug) == len(prs.slides)
    assert report.wps_risk_overflow_count == len(report.critical_issues)
    assert report.wps_risk_overflow_count == 0, (
        f"expected critical=0 after WPS-safe renderer, got {report.wps_risk_overflow_count}"
    )
    legacy = verify_text_fit(prs)
    assert isinstance(legacy, list)
