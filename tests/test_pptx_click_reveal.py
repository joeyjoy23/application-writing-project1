"""Tests for pptx_click_reveal."""

import re
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from scripts.generate_classroom_pptx import SlideBuilder
from scripts.generate_classroom_pptx_v2 import render_v2_deck
from scripts.pptx_click_reveal import apply_click_reveal, count_click_reveal_stats

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_ANIM_RE = re.compile(r"^anim_\d{3}$")


def _anim_names_from_pptx(path: Path) -> list[str]:
    names: list[str] = []
    with zipfile.ZipFile(path, "r") as zf:
        for name in sorted(zf.namelist()):
            if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
                continue
            root = ET.fromstring(zf.read(name))
            for cnv in root.iter(f"{{{_P_NS}}}cNvPr"):
                n = cnv.get("name", "")
                if _ANIM_RE.match(n):
                    names.append(n)
    return names


def _make_anim_pptx(path: Path) -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for i in range(1, 4):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1), Inches(i), Inches(3), Inches(0.5)
        )
        box.name = f"anim_{i:03d}"
    prs.save(path)


def test_apply_click_reveal_injects_timing(tmp_path):
    pptx = tmp_path / "test.pptx"
    _make_anim_pptx(pptx)
    ok = apply_click_reveal(pptx)
    assert ok is True
    with zipfile.ZipFile(pptx, "r") as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml")
    root = ET.fromstring(slide_xml)
    timing = root.find(f"{{{_P_NS}}}timing")
    assert timing is not None
    assert b"p:timing" in slide_xml or timing.tag.endswith("timing")
    assert b"attrNameLst" in slide_xml or b"attrName" in slide_xml
    assert b"presetID=\"10\"" in slide_xml
    assert b"animEffect" in slide_xml
    assert b"bldLst" in slide_xml


def test_apply_click_reveal_table_shapes(tmp_path):
    pptx = tmp_path / "table_anim.pptx"
    prs = Presentation()
    builder = SlideBuilder(prs)
    builder.vocab_table_slide(
        "话题词块 · 测试",
        "必备级",
        [{"english": "cracked heart", "example": "The cracked heart symbolizes vulnerability."}],
        ["english", "example"],
    )
    prs.save(pptx)
    ok = apply_click_reveal(pptx)
    assert ok is True
    with zipfile.ZipFile(pptx, "r") as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml")
    root = ET.fromstring(slide_xml)
    timing = root.find(f"{{{_P_NS}}}timing")
    assert timing is not None


def test_apply_click_reveal_keeps_anim_shapes_visible(tmp_path):
    pptx = tmp_path / "test.pptx"
    _make_anim_pptx(pptx)
    apply_click_reveal(pptx)
    with zipfile.ZipFile(pptx, "r") as zf:
        slide_xml = zf.read("ppt/slides/slide1.xml")
    root = ET.fromstring(slide_xml)
    hidden = [
        c.get("name")
        for host in root.iter(f"{{{_P_NS}}}sp")
        for c in host.iter(f"{{{_P_NS}}}cNvPr")
        if c.get("hidden") == "1"
    ]
    assert hidden == []


def test_apply_click_reveal_no_anim_shapes(tmp_path):
    pptx = tmp_path / "plain.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(pptx)
    ok = apply_click_reveal(pptx)
    assert ok is True


def test_v2_content_slide_click_reveal(tmp_path):
    out = tmp_path / "v2_anim.pptx"
    slides = [
        {
            "type": "content",
            "title": "审题 · 测试页",
            "bullets": ["要点一：读懂题目", "要点二：明确对象", "↓"],
        },
        {
            "type": "content",
            "title": "抓重点",
            "bullets": ["★ 高分关键句"],
            "callout": "本题核心：心理健康与海报设计",
        },
    ]
    render_v2_deck(slides, out)
    names_before = _anim_names_from_pptx(out)
    assert len(names_before) >= 4
    assert all(_ANIM_RE.match(n) for n in names_before)

    ok = apply_click_reveal(out)
    assert ok is True
    timing_slides, anim_count = count_click_reveal_stats(out)
    assert timing_slides >= 2
    assert anim_count >= 4

    with zipfile.ZipFile(out, "r") as zf:
        slide1 = ET.fromstring(zf.read("ppt/slides/slide1.xml"))
        assert slide1.find(f"{{{_P_NS}}}timing") is not None
    hidden = [
        c.get("name")
        for host in slide1.iter(f"{{{_P_NS}}}sp")
        for c in host.iter(f"{{{_P_NS}}}cNvPr")
        if c.get("hidden") == "1" and _ANIM_RE.match(c.get("name", ""))
    ]
    assert hidden == []
