"""Unified PPT card primitives — white surface, mint accent bar, subtle shadow."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from styles.design_tokens import Theme as G

CARD_RADIUS_ADJ = 0.10
ACCENT_BAR_IN = 0.042
CARD_PAD_TOP = 0.10
CARD_PAD_LEFT = 0.16
CARD_INNER_LEFT = 0.20


def apply_subtle_shadow(shape) -> None:
    """Very light drop shadow (Notion-like, not decorative)."""
    sp_pr = shape._sp.spPr
    for old in sp_pr.findall(qn("a:effectLst")):
        sp_pr.remove(old)
    effect_lst = OxmlElement("a:effectLst")
    outer = OxmlElement("a:outerShdw")
    outer.set("blurRad", "25400")
    outer.set("dist", "12700")
    outer.set("dir", "5400000")
    outer.set("algn", "tl")
    outer.set("rotWithShape", "0")
    srgb = OxmlElement("a:srgbClr")
    srgb.set("val", "1F2937")
    alpha = OxmlElement("a:alpha")
    alpha.set("val", "10000")
    srgb.append(alpha)
    outer.append(srgb)
    effect_lst.append(outer)
    sp_pr.append(effect_lst)


def add_unified_card(
    slide,
    left,
    top,
    width,
    height,
    *,
    accent: RGBColor | None = None,
    warning: bool = False,
    fill: RGBColor | None = None,
):
    """White card + 3px mint (or amber) left bar + light border + shadow."""
    accent = accent or G.card_accent(warning=warning)
    fill = fill or G.SURFACE
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = G.BORDER
    card.line.width = Pt(1)
    if card.adjustments:
        card.adjustments[0] = CARD_RADIUS_ADJ
    apply_subtle_shadow(card)

    pad_v = Inches(0.08)
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left,
        top + pad_v,
        Inches(ACCENT_BAR_IN),
        height - pad_v * 2,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    return card


def inner_box_left(left) -> float:
    return float(left.inches if hasattr(left, "inches") else left) + CARD_INNER_LEFT


def inner_box_width(width) -> float:
    w = float(width.inches if hasattr(width, "inches") else width)
    return max(1.0, w - CARD_INNER_LEFT - 0.14)
